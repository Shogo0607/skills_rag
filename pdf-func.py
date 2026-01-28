from __future__ import annotations

import io
import logging
import os
import mimetypes
import textwrap
from datetime import datetime
import threading
from pathlib import Path
from typing import Any, Iterator, List, Literal, Sequence, TypedDict
import shutil
import re
import subprocess
import platform

from fastapi import BackgroundTasks, FastAPI, HTTPException, Query, UploadFile, File, Form, Request
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, StreamingResponse
from pydantic import BaseModel
from openai import OpenAI, AzureOpenAI
import secrets
import json
import urllib.parse
import urllib.request
from dotenv import load_dotenv
import hmac
import hashlib
import base64
import time
import pdfplumber
from langchain_core.language_models.chat_models import BaseChatModel
from langchain_core.messages import BaseMessage, HumanMessage, SystemMessage
from langchain_openai import AzureChatOpenAI, ChatOpenAI
from PIL import Image, ImageDraw, ImageFont
try:
    # OpenAI Agents SDK (prompt-driven MCP integration)
    from agents import Agent, Runner  # type: ignore
    from agents.model_settings import ModelSettings  # type: ignore
    from agents.mcp import MCPServerStreamableHttp  # type: ignore
except Exception:
    Agent = None  # type: ignore[assignment]
    Runner = None  # type: ignore[assignment]
    ModelSettings = None  # type: ignore[assignment]
    MCPServerStreamableHttp = None  # type: ignore[assignment]

try:
    from pptx import Presentation as PptxPresentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except Exception:  # pragma: no cover - optional at import time; endpoint will validate
    PptxPresentation = None  # type: ignore[assignment]
    MSO_SHAPE_TYPE = None  # type: ignore[assignment]

try:
    from docx import Document as DocxDocument
    from docx.oxml.table import CT_Tbl
    from docx.oxml.text.paragraph import CT_P
    from docx.table import Table as DocxTable
    from docx.text.paragraph import Paragraph as DocxParagraph
except Exception:  # pragma: no cover - optional at import time; endpoint will validate
    DocxDocument = None  # type: ignore[assignment]
    DocxTable = None  # type: ignore[assignment]
    DocxParagraph = None  # type: ignore[assignment]
    CT_Tbl = None  # type: ignore[assignment]
    CT_P = None  # type: ignore[assignment]

WORD_NAMESPACE = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
W_P = f"{{{WORD_NAMESPACE}}}p"
W_T = f"{{{WORD_NAMESPACE}}}t"

# NOTE: We no longer load API keys from .env; keys are stored per-user in the profile store.
# Keep dotenv to optionally load non-sensitive values like SHARE_SECRET if present.
load_dotenv(".env")

try:
    import openpyxl  # type: ignore
except Exception:  # pragma: no cover - optional at import time; endpoint will validate
    openpyxl = None


class Entry(TypedDict):
    name: str
    path: str
    type: Literal["file", "directory"]
    size: int | None
    mtime: float | None
    isLink: bool


app = FastAPI(title="Directory Explorer API", version="0.1.0")


# Configure base directory (database folder under workspace)
# Resolve workspace root relative to this file: backend/ -> workspace root
WORKSPACE_ROOT = Path(__file__).resolve().parent.parent
BASE_DIR = (WORKSPACE_ROOT / "database").resolve()

BASE_DIR.mkdir(parents=True, exist_ok=True)

logger = logging.getLogger("directory_rag.pdf_extraction")

# -------------------- In-memory upload/progress tracking --------------------
_PROGRESS_LOCK = threading.Lock()
_UPLOAD_PROGRESS: dict[str, dict] = {}

def _progress_key(username: str, rel_path: str) -> str:
    key_path = rel_path.replace("\\", "/").strip("/")
    return f"{username}:{key_path}"

def _set_progress(username: str, rel_path: str, **updates: Any) -> None:
    key = _progress_key(username, rel_path)
    with _PROGRESS_LOCK:
        data = _UPLOAD_PROGRESS.get(key) or {}
        data.update(updates)
        if "percent" in data:
            try:
                p = float(data["percent"]) if data["percent"] is not None else 0.0
            except Exception:
                p = 0.0
            data["percent"] = max(0.0, min(1.0, p))
        _UPLOAD_PROGRESS[key] = data

def _get_progress(username: str, rel_path: str) -> dict:
    key = _progress_key(username, rel_path)
    with _PROGRESS_LOCK:
        return dict(_UPLOAD_PROGRESS.get(key) or {})

def _clear_progress(username: str, rel_path: str) -> None:
    key = _progress_key(username, rel_path)
    with _PROGRESS_LOCK:
        _UPLOAD_PROGRESS.pop(key, None)

# ---- Public Agent registry (simple file-backed KV in database/_users/agents.json) ----
AGENTS_FILE = (BASE_DIR / "_users" / "agents.json")

# ---- Public Shares registry (simple file-backed KV in database/_users/shares.json) ----
SHARES_FILE = (BASE_DIR / "_users" / "shares.json")

def _ensure_agents_file() -> None:
    try:
        (BASE_DIR / "_users").mkdir(parents=True, exist_ok=True)
        if not AGENTS_FILE.exists():
            AGENTS_FILE.write_text("{}", encoding="utf-8")
    except Exception:
        raise HTTPException(status_code=500, detail="Failed to initialize agents store")

def _load_agents() -> dict:
    _ensure_agents_file()
    try:
        return json.loads(AGENTS_FILE.read_text(encoding="utf-8"))
    except Exception:
        return {}

def _save_agents(data: dict) -> None:
    _ensure_agents_file()
    try:
        AGENTS_FILE.write_text(json.dumps(data, ensure_ascii=False, indent=2), encoding="utf-8")
    except Exception:
        raise HTTPException(status_code=500, detail="Failed to save agents store")

def _ensure_shares_file() -> None:
    try:
        (BASE_DIR / "_users").mkdir(parents=True, exist_ok=True)
        if not SHARES_FILE.exists():
            SHARES_FILE.write_text("{}", encoding="utf-8")
    except Exception:
        raise HTTPException(status_code=500, detail="Failed to initialize shares store")

def _load_shares() -> dict:
    _ensure_shares_file()
    try:
        return json.loads(SHARES_FILE.read_text(encoding="utf-8"))
    except Exception:
        return {}

def _save_shares(data: dict) -> None:
    _ensure_shares_file()
    try:
        SHARES_FILE.write_text(json.dumps(data, ensure_ascii=False, indent=2), encoding="utf-8")
    except Exception:
        raise HTTPException(status_code=500, detail="Failed to save shares store")

def _build_summary_system_prompt(document_label_jp: str, unit_label_jp: str) -> str:
    return textwrap.dedent(
        f"""
        あなたは{document_label_jp}の各{unit_label_jp}を要約する専門家です。
        重要なトピック、数字、箇条書き項目などを抽出し、日本語で簡潔にまとめてください。
        出力は最大3つの箇条書きで、冗長な表現は避けてください。
        """
    ).strip()


def _build_category_system_prompt(document_label_jp: str, unit_label_jp: str) -> str:
    return textwrap.dedent(
        f"""
        あなたは{document_label_jp}の分析と分類を担当するアナリストです。
        入力として{unit_label_jp}番号とその簡潔な要約が与えられます。
        それらを意味のあるカテゴリに分け、カテゴリ名と説明を日本語で作成してください。
        各{unit_label_jp}は必ずいずれか1つのカテゴリに割り当ててください。
        出力は以下のJSON形式のみを使用してください。
        {{
          "categories": [
            {{ "name": "カテゴリ名", "description": "カテゴリの説明", "pages": [{unit_label_jp}番号, ...] }}
          ]
        }}
        """
    ).strip()


PDF_EXTRACTION_SYSTEM_PROMPT = textwrap.dedent(
    """
    You are a meticulous document transcription expert.
    Transcribe every piece of textual information from provided page images,
    including headers, footers, captions, annotations inside figures, callouts,
    and all table contents. Preserve the logical structure using Markdown headings,
    bullet lists, and Markdown tables. Do not omit numeric values or labels.
    """
).strip()

PDF_EXTRACTION_PAGE_INSTRUCTIONS = textwrap.dedent(
    """
    Extract everything visible on this page image. Preserve reading order when possible.
    Render tables as Markdown tables and include unit labels. If a figure or diagram
    contains text, reproduce it verbatim and provide a short description when helpful.
    Maintain Japanese text as-is. Respond in Markdown only.
    """
).strip()

PAGE_SUMMARY_SYSTEM_PROMPT = _build_summary_system_prompt("PDF文書", "ページ")

CATEGORY_SYSTEM_PROMPT = _build_category_system_prompt("PDF文書", "ページ")

PPTX_EXTRACTION_SYSTEM_PROMPT = textwrap.dedent(
    """
    You are a meticulous transcription specialist for PowerPoint slides.
    Capture every textual element visible on the slide image as well as meaningful descriptions
    of visual elements such as charts, diagrams, icons, highlighted regions, and annotations.
    Structure the output with Markdown headings, bullet lists, and Markdown tables whenever appropriate.
    Do not omit labels, numeric values, or descriptive captions embedded inside visuals.
    """
).strip()

class DocumentProcessingPrompts(TypedDict):
    doc_label_jp: str
    unit_label_en: str
    unit_label_jp: str
    unit_prefix: str
    extraction_system_prompt: str
    extraction_instruction: str
    summary_system_prompt: str
    category_system_prompt: str


EMU_PER_INCH = 914400
PPTX_RENDER_DPI = 150
_FONT_CACHE: dict[int, ImageFont.ImageFont] = {}
_FONT_CANDIDATES = [
    "NotoSansCJKjp-Regular.otf",
    "NotoSansCJK-Regular.ttc",
    "NotoSansJP-Regular.otf",
    "SourceHanSans-Regular.otf",
    "Hiragino Sans W3.ttc",
    "HiraginoSans-W3.ttc",
    "YuGothR.ttc",
    "MSGothic.ttc",
    "Arial Unicode.ttf",
    "TakaoPGothic.ttf",
    "ipaexg.ttf",
    "DejaVuSans.ttf",
]
_FONT_DIRECTORIES = [
    Path("/System/Library/Fonts"),
    Path("/System/Library/Fonts/Supplemental"),
    Path("/Library/Fonts"),
    Path("~/Library/Fonts").expanduser(),
    Path("/usr/share/fonts"),
    Path("/usr/local/share/fonts"),
]

PPTX_EXTRACTION_SLIDE_INSTRUCTIONS = textwrap.dedent(
    """
    このスライドに含まれるテキストと視覚要素をすべて抽出してください。
    図表・画像・アイコンなどに書かれている文字や強調表現も正確に読み取り、必要に応じて短い説明を付けてください。
    補助テキストが与えられている場合は照合しつつ、画像から読み取れる情報を優先してMarkdown形式で出力してください。
    """
).strip()


def _ensure_output_dir(pdf_path: Path) -> Path:
    parent = pdf_path.parent
    if parent.name == pdf_path.stem:
        output_dir = parent
    else:
        output_dir = parent / pdf_path.stem
        output_dir.mkdir(parents=True, exist_ok=True)
    return output_dir


def _cleanup_previous_outputs(output_dir: Path, preserve: set[Path]) -> None:
    preserved: set[Path] = set()
    for p in preserve:
        try:
            preserved.add(p.resolve())
        except Exception:
            preserved.add(p)
    try:
        for child in list(output_dir.iterdir()):
            try:
                resolved_child = child.resolve()
            except Exception:
                resolved_child = child
            if resolved_child in preserved:
                continue
            try:
                if child.is_dir():
                    shutil.rmtree(child)
                else:
                    child.unlink()
            except Exception:
                logger.exception("Failed to remove %s within %s", child, output_dir)
    except Exception:
        logger.exception("Failed to clean existing outputs in %s", output_dir)


def _write_summary(output_dir: Path, pdf_name: str, message: str) -> None:
    summary_path = output_dir / "summary.md"
    try:
        summary_content = f"# Extracted Content for {pdf_name}\n\n{message.strip()}\n"
        summary_path.write_text(summary_content, encoding="utf-8")
    except Exception:
        logger.exception("Failed to write summary for %s", pdf_name)


def _pil_image_to_data_url(image: Image.Image) -> str:
    buffer = io.BytesIO()
    image.save(buffer, format="PNG")
    encoded = base64.b64encode(buffer.getvalue()).decode("ascii")
    return f"data:image/png;base64,{encoded}"


def _message_content_to_text(message: BaseMessage) -> str:
    content = message.content
    if isinstance(content, str):
        return content
    if isinstance(content, list):
        parts: list[str] = []
        for block in content:
            if isinstance(block, dict):
                text = block.get("text")
                if text:
                    parts.append(str(text))
        return "\n".join(parts)
    return str(content)


def _iter_docx_blocks(document: DocxDocument) -> Iterator[DocxParagraph | DocxTable]:
    if DocxParagraph is None or DocxTable is None or CT_P is None or CT_Tbl is None:
        return
    body = document.element.body  # type: ignore[attr-defined]
    for child in body.iterchildren():
        if CT_P is not None and isinstance(child, CT_P):
            yield DocxParagraph(child, document)  # type: ignore[arg-type]
        elif CT_Tbl is not None and isinstance(child, CT_Tbl):
            yield DocxTable(child, document)  # type: ignore[arg-type]


def _collect_docx_paragraph_texts(element: Any) -> list[str]:
    if element is None or not hasattr(element, "iter"):
        return []
    paragraphs: list[str] = []
    for node in element.iter():
        tag = getattr(node, "tag", "")
        if tag == W_P:
            runs: list[str] = []
            for text_node in node.iter():
                if getattr(text_node, "tag", "") == W_T and text_node.text:
                    runs.append(text_node.text)
            para_text = "".join(runs).replace("\xa0", " ").strip()
            if para_text:
                paragraphs.append(para_text)
    if not paragraphs:
        fallback: list[str] = []
        for text_node in element.iter():
            if getattr(text_node, "tag", "") == W_T and text_node.text:
                value = text_node.text.replace("\xa0", " ").strip()
                if value:
                    fallback.append(value)
        if fallback:
            paragraphs.append("".join(fallback))
    return paragraphs


def _docx_paragraph_to_text(paragraph: DocxParagraph) -> str:
    text = paragraph.text.replace("\xa0", " ").strip()
    if text:
        return text
    oxml = getattr(paragraph, "_p", None)
    if oxml is None:
        return ""
    return "\n".join(_collect_docx_paragraph_texts(oxml)).strip()


def _docx_cell_to_text(cell: Any) -> str:
    value = str(getattr(cell, "text", "") or "").replace("\xa0", " ").strip()
    if value:
        return value
    oxml = getattr(cell, "_tc", None)
    if oxml is None:
        return ""
    return "\n".join(_collect_docx_paragraph_texts(oxml)).strip()


def _get_font(size: int) -> ImageFont.ImageFont:
    cached = _FONT_CACHE.get(size)
    if cached is not None:
        return cached
    for name in _FONT_CANDIDATES:
        try:
            font = ImageFont.truetype(name, size=size)
            _FONT_CACHE[size] = font
            return font
        except Exception:
            pass
        for directory in _FONT_DIRECTORIES:
            candidate = directory / name
            if not candidate.exists():
                continue
            try:
                font = ImageFont.truetype(str(candidate), size=size)
                _FONT_CACHE[size] = font
                return font
            except Exception:
                continue
    font = ImageFont.load_default()
    _FONT_CACHE[size] = font
    return font


def _text_width(draw: ImageDraw.ImageDraw, text: str, font: ImageFont.ImageFont) -> int:
    try:
        return int(draw.textlength(text, font=font))
    except AttributeError:
        bbox = draw.textbbox((0, 0), text, font=font)
        if bbox is None:
            return 0
        return int(bbox[2] - bbox[0])


def _wrap_text_to_width(text: str, font: ImageFont.ImageFont, max_width: int, draw: ImageDraw.ImageDraw) -> list[str]:
    if max_width <= 0:
        return [text]
    lines: list[str] = []
    for raw_paragraph in text.replace("\r", "").split("\n"):
        paragraph = raw_paragraph.strip()
        if not paragraph:
            lines.append("")
            continue
        words = paragraph.split()
        if not words:  # likely CJK text without spaces
            current = ""
            for char in paragraph:
                candidate = current + char
                if not current or _text_width(draw, candidate, font) <= max_width:
                    current = candidate
                else:
                    lines.append(current)
                    current = char
            if current:
                lines.append(current)
            continue
        current = ""
        for word in words:
            candidate = f"{current} {word}".strip() if current else word
            if _text_width(draw, candidate, font) <= max_width or not current:
                current = candidate
            else:
                if current:
                    lines.append(current)
                current = word
        if current:
            lines.append(current)
    return lines


def _draw_text_block(
    draw: ImageDraw.ImageDraw,
    text: str,
    left: int,
    top: int,
    width: int,
    height: int,
    *,
    font_size: int | None = None,
) -> None:
    cleaned = text.strip()
    if not cleaned or width <= 0 or height <= 0:
        return
    effective_font_size = font_size or max(16, min(48, height // 4 or 16))
    font = _get_font(effective_font_size)
    padding = max(4, min(width, height) // 30)
    available_width = max(1, width - padding * 2)
    lines = _wrap_text_to_width(cleaned, font, available_width, draw)
    try:
        ascent, descent = font.getmetrics()
        line_height = ascent + descent + 2
    except Exception:
        line_height = effective_font_size + 4
    y = top + padding
    for line in lines:
        if y + line_height > top + height - padding:
            break
        draw.text((left + padding, y), line, fill="black", font=font)
        y += line_height


def _emu_to_px(value: int | float, dpi: int = PPTX_RENDER_DPI) -> int:
    try:
        numeric = float(value)
    except Exception:
        numeric = 0.0
    inches = numeric / EMU_PER_INCH
    pixels = inches * dpi
    return max(1, int(round(pixels)))


def _render_pptx_shape(
    image: Image.Image,
    draw: ImageDraw.ImageDraw,
    shape: Any,
    dpi: int,
    offset_left: int = 0,
    offset_top: int = 0,
) -> None:
    try:
        shape_type = shape.shape_type if MSO_SHAPE_TYPE is not None else None
    except Exception:
        shape_type = None
    try:
        left = offset_left + _emu_to_px(int(shape.left), dpi)
        top = offset_top + _emu_to_px(int(shape.top), dpi)
        width = _emu_to_px(int(shape.width), dpi)
        height = _emu_to_px(int(shape.height), dpi)
    except Exception:
        left = offset_left
        top = offset_top
        width = 0
        height = 0

    if MSO_SHAPE_TYPE is not None and shape_type == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:
            _render_pptx_shape(image, draw, child, dpi, left, top)
        return

    if width <= 0 or height <= 0:
        return

    if MSO_SHAPE_TYPE is not None and shape_type == MSO_SHAPE_TYPE.PICTURE:
        try:
            picture = Image.open(io.BytesIO(shape.image.blob)).convert("RGBA")
            resized = picture.resize((max(1, width), max(1, height)), Image.LANCZOS)
            image.paste(resized, (left, top), resized if resized.mode == "RGBA" else None)
            return
        except Exception:
            _draw_text_block(draw, "[画像]", left, top, width, height)
            return

    if getattr(shape, "has_table", False):
        rows: list[str] = []
        try:
            table = shape.table
        except Exception:
            table = None
        if table is not None:
            for row in table.rows:
                cells = []
                for cell in row.cells:
                    value = (cell.text or "").strip()
                    if len(value) > 120:
                        value = value[:117] + "..."
                    cells.append(value)
                rows.append(" | ".join(cells).strip())
        table_text = "\n".join(filter(None, rows))
        if table_text:
            _draw_text_block(draw, table_text, left, top, width, height, font_size=max(12, min(28, height // 6 or 12)))
        return

    if getattr(shape, "has_text_frame", False):
        paragraphs: list[str] = []
        try:
            for paragraph in shape.text_frame.paragraphs:
                para_text = "".join(run.text for run in paragraph.runs) if getattr(paragraph, "runs", None) else paragraph.text
                if para_text:
                    paragraphs.append(para_text.strip())
        except Exception:
            text_val = getattr(shape, "text", "")
            if isinstance(text_val, str):
                paragraphs.append(text_val.strip())
        text_content = "\n".join(filter(None, paragraphs))
        if text_content:
            _draw_text_block(draw, text_content, left, top, width, height)
        return

    if MSO_SHAPE_TYPE is not None and shape_type == MSO_SHAPE_TYPE.CHART:
        title = ""
        chart = getattr(shape, "chart", None)
        if chart is not None:
            try:
                if chart.has_title and chart.chart_title.has_text_frame:
                    title = chart.chart_title.text_frame.text.strip()
            except Exception:
                title = ""
        label = "[チャート]" + (f" {title}" if title else "")
        _draw_text_block(draw, label, left, top, width, height)
        return

    if MSO_SHAPE_TYPE is not None and shape_type == MSO_SHAPE_TYPE.SMART_ART:
        _draw_text_block(draw, "[SmartArt]", left, top, width, height)
        return

    if MSO_SHAPE_TYPE is not None and shape_type == MSO_SHAPE_TYPE.MEDIA:
        _draw_text_block(draw, "[メディア]", left, top, width, height)
        return

    text_value = getattr(shape, "text", "")
    if isinstance(text_value, str) and text_value.strip():
        _draw_text_block(draw, text_value.strip(), left, top, width, height)


def _render_pptx_slide(slide: Any, presentation: Any, dpi: int = PPTX_RENDER_DPI) -> Image.Image:
    width_px = max(640, _emu_to_px(int(getattr(presentation, "slide_width", 0) or 0), dpi))
    height_px = max(360, _emu_to_px(int(getattr(presentation, "slide_height", 0) or 0), dpi))
    image = Image.new("RGB", (width_px, height_px), "white")
    draw = ImageDraw.Draw(image)
    for shape in getattr(slide, "shapes", []):
        try:
            _render_pptx_shape(image, draw, shape, dpi)
        except Exception:
            logger.exception("Failed to render shape on PPTX slide")
    return image


def _describe_pptx_shape(shape: Any, depth: int = 0) -> list[str]:
    prefix = "  " * depth + "- "
    try:
        shape_type = shape.shape_type if MSO_SHAPE_TYPE is not None else None
    except Exception:
        shape_type = None

    name = getattr(shape, "name", "")
    name_suffix = f"（{name.strip()}）" if isinstance(name, str) and name.strip() else ""

    def truncate(value: str, limit: int = 240) -> str:
        value = value.strip()
        if len(value) > limit:
            return value[: limit - 3] + "..."
        return value

    lines: list[str] = []
    if MSO_SHAPE_TYPE is not None and shape_type == MSO_SHAPE_TYPE.GROUP:
        lines.append(f"{prefix}グループ{name_suffix}".rstrip(": "))
        for child in shape.shapes:
            lines.extend(_describe_pptx_shape(child, depth + 1))
        return lines

    if getattr(shape, "has_table", False):
        rows: list[str] = []
        try:
            for row in shape.table.rows:
                cells = [truncate(cell.text or "") for cell in row.cells]
                rows.append(" | ".join(filter(None, cells)))
        except Exception:
            rows = []
        payload = "; ".join(filter(None, rows))
        if payload:
            lines.append(f"{prefix}表{name_suffix}: {payload}")
        else:
            lines.append(f"{prefix}表{name_suffix}")
        return lines

    if getattr(shape, "has_text_frame", False):
        fragments: list[str] = []
        try:
            for paragraph in shape.text_frame.paragraphs:
                text_val = "".join(run.text for run in paragraph.runs) if getattr(paragraph, "runs", None) else paragraph.text
                if text_val:
                    fragments.append(truncate(text_val))
        except Exception:
            text_val = getattr(shape, "text", "")
            if isinstance(text_val, str) and text_val.strip():
                fragments.append(truncate(text_val))
        payload = " / ".join(filter(None, fragments))
        if payload:
            lines.append(f"{prefix}テキスト{name_suffix}: {payload}")
        return lines

    if MSO_SHAPE_TYPE is not None and shape_type == MSO_SHAPE_TYPE.PICTURE:
        alt = getattr(shape, "alternative_text", "") or getattr(shape, "description", "")
        payload = truncate(alt) if isinstance(alt, str) else ""
        lines.append(f"{prefix}画像{name_suffix}" + (f": {payload}" if payload else ""))
        return lines

    if MSO_SHAPE_TYPE is not None and shape_type == MSO_SHAPE_TYPE.CHART:
        title = ""
        chart = getattr(shape, "chart", None)
        if chart is not None:
            try:
                if chart.has_title and chart.chart_title.has_text_frame:
                    title = chart.chart_title.text_frame.text.strip()
            except Exception:
                title = ""
        payload = truncate(title) if title else ""
        lines.append(f"{prefix}チャート{name_suffix}" + (f": {payload}" if payload else ""))
        return lines

    if MSO_SHAPE_TYPE is not None and shape_type == MSO_SHAPE_TYPE.SMART_ART:
        lines.append(f"{prefix}SmartArt{name_suffix}")
        return lines

    if MSO_SHAPE_TYPE is not None and shape_type == MSO_SHAPE_TYPE.MEDIA:
        lines.append(f"{prefix}メディア{name_suffix}")
        return lines

    text_val = getattr(shape, "text", "")
    if isinstance(text_val, str) and text_val.strip():
        lines.append(f"{prefix}テキスト{name_suffix}: {truncate(text_val)}")
    return lines


def _extract_pptx_slide_context(slide: Any) -> str:
    parts: list[str] = []
    for shape in getattr(slide, "shapes", []):
        try:
            parts.extend(_describe_pptx_shape(shape))
        except Exception:
            logger.exception("Failed to describe shape within PPTX slide")
    notes: list[str] = []
    try:
        if getattr(slide, "has_notes_slide", False) and slide.notes_slide is not None:
            for shape in slide.notes_slide.shapes:
                if getattr(shape, "has_text_frame", False):
                    value = (shape.text or "").strip()
                    if value:
                        notes.append(value)
    except Exception:
        logger.exception("Failed to extract notes from PPTX slide")
    if notes:
        parts.append(f"- 発表者ノート: {' / '.join(notes)}")
    return "\n".join(part for part in parts if part)


def _build_messages_for_unit(
    unit_number: int,
    image_data_url: str,
    prompts: DocumentProcessingPrompts,
    *,
    extra_text: str | None = None,
) -> Sequence[BaseMessage]:
    unit_label = prompts["unit_label_en"]
    text_instruction = f"{unit_label} {unit_number:03}: {prompts['extraction_instruction']}"
    if extra_text:
        text_instruction = f"{text_instruction}\n\n補助テキスト:\n{extra_text.strip()}"
    return [
        SystemMessage(content=prompts["extraction_system_prompt"]),
        HumanMessage(
            content=[
                {"type": "text", "text": text_instruction},
                {"type": "image_url", "image_url": {"url": image_data_url, "detail": "high"}},
            ]
        ),
    ]


def _build_summary_messages_for_unit(
    unit_number: int,
    unit_markdown: str,
    prompts: DocumentProcessingPrompts,
) -> Sequence[BaseMessage]:
    snippet = unit_markdown.strip()
    if len(snippet) > 6000:
        snippet = snippet[:6000] + "\n...[truncated]"
    instruction = textwrap.dedent(
        f"""
        以下は{prompts['doc_label_jp']}の{prompts['unit_label_jp']} {unit_number:03}の内容です。重要な論点や数字、表の内容を最大3つの箇条書きで要約してください。
        各箇条書きは40文字以内を目安とし、日本語で記載してください。
        """
    ).strip()
    return [
        SystemMessage(content=prompts["summary_system_prompt"]),
        HumanMessage(
            content=[
                {"type": "text", "text": f"{instruction}\n\n```markdown\n{snippet}\n```"}
            ]
        ),
    ]


def _summarize_units(model: BaseChatModel, unit_entries: list[dict], prompts: DocumentProcessingPrompts) -> dict[int, str]:
    if not unit_entries:
        return {}
    messages = [
        _build_summary_messages_for_unit(entry["number"], entry["text"], prompts) for entry in unit_entries
    ]
    responses = model.batch(messages, config={"max_concurrency": 4})
    summaries: dict[int, str] = {}
    for entry, response in zip(unit_entries, responses):
        summary = _message_content_to_text(response).strip()
        summaries[entry["number"]] = summary or "(要約なし)"
    return summaries


def _categorize_units(
    model: BaseChatModel,
    unit_summaries: dict[int, str],
    prompts: DocumentProcessingPrompts,
) -> list[dict]:
    if not unit_summaries:
        return []
    summary_lines = []
    unit_label = prompts["unit_label_en"]
    for number in sorted(unit_summaries.keys()):
        summary = unit_summaries[number]
        snippet = summary if len(summary) <= 1200 else summary[:1200] + "..."
        summary_lines.append(f"{unit_label} {number:03}: {snippet}")
    joined = "\n".join(summary_lines)
    messages = [
        SystemMessage(content=prompts["category_system_prompt"]),
        HumanMessage(content=[{"type": "text", "text": joined}]),
    ]
    response = model.invoke(messages)
    raw_text = _message_content_to_text(response)
    data = None
    try:
        data = json.loads(raw_text)
    except Exception:
        match = re.search(r"\{.*\}", raw_text, re.DOTALL)
        if match:
            try:
                data = json.loads(match.group(0))
            except Exception:
                data = None

    if not isinstance(data, dict):
        return [
            {
                "name": f"{unit_label} {num:03}",
                "description": "自動分類に失敗したため個別カテゴリとして保持します。",
                "pages": [num],
            }
            for num in sorted(unit_summaries.keys())
        ]

    categories: list[dict] = []
    for idx, cat in enumerate(data.get("categories") or [], start=1):
        pages = []
        for p in cat.get("pages", []):
            try:
                pn = int(p)
            except Exception:
                continue
            pages.append(pn)
        pages = sorted(set(pages))
        if not pages:
            continue
        name = str(cat.get("name") or f"Category {idx}")
        description = str(cat.get("description") or "")
        categories.append({"name": name, "description": description, "pages": pages})

    if not categories:
        return [{"name": "Uncategorized", "description": "", "pages": sorted(unit_summaries.keys())}]

    # Ensure every page appears at least once; add leftovers to Uncategorized
    assigned = {p for cat in categories for p in cat["pages"]}
    leftovers = sorted(set(unit_summaries.keys()) - assigned)
    if leftovers:
        categories.append({"name": "Uncategorized", "description": "自動分類に割り当てられなかったページです。", "pages": leftovers})

    return categories


def _slugify_category_name(name: str, fallback: str) -> str:
    cleaned = re.sub(r"[\\/:*?\"<>|]", "_", name.strip())
    cleaned = cleaned.strip(" ._")
    if not cleaned:
        cleaned = fallback
    if len(cleaned) > 40:
        cleaned = cleaned[:40].rstrip(" ._") or cleaned[:40]
    return cleaned or fallback


def _ensure_unique_folder_name(base_name: str, used: set[str]) -> str:
    candidate = base_name
    counter = 1
    while candidate in used:
        candidate = f"{base_name}_{counter:02}"
        counter += 1
    used.add(candidate)
    return candidate


def _organize_units_by_category(
    output_dir: Path,
    unit_entries: list[dict],
    categories: list[dict],
    unit_summaries: dict[int, str],
    prompts: DocumentProcessingPrompts,
) -> str:
    unit_lookup = {entry["number"]: entry for entry in unit_entries}
    used_folder_names: set[str] = set()
    assigned: set[int] = set()
    category_details: list[dict] = []

    for idx, category in enumerate(categories, start=1):
        units = [p for p in category.get("pages", []) if p in unit_lookup and p not in assigned]
        if not units:
            continue
        base = _slugify_category_name(category.get("name", f"Category {idx}"), f"category_{idx:02}")
        folder_name = _ensure_unique_folder_name(base, used_folder_names)
        folder_path = output_dir / folder_name
        folder_path.mkdir(parents=True, exist_ok=True)
        moved_units = []
        for p in units:
            entry = unit_lookup[p]
            source = entry["md_path"]
            dest = folder_path / source.name
            try:
                if dest.exists():
                    dest.unlink()
                shutil.move(str(source), str(dest))
            except Exception:
                logger.exception("Failed to move %s to %s", source, dest)
                continue
            entry["md_path"] = dest
            assigned.add(p)
            moved_units.append(
                {
                    "number": p,
                    "summary": unit_summaries.get(p, ""),
                    "relative_path": dest.relative_to(output_dir).as_posix(),
                }
            )
        if moved_units:
            category_details.append(
                {
                    "name": category.get("name") or folder_name,
                    "description": category.get("description", ""),
                    "folder": folder_name,
                    "pages": moved_units,
                }
            )

    leftovers = [entry for entry in unit_entries if entry["number"] not in assigned]
    if leftovers:
        folder_name = _ensure_unique_folder_name("Uncategorized", used_folder_names)
        folder_path = output_dir / folder_name
        folder_path.mkdir(parents=True, exist_ok=True)
        moved_units = []
        for entry in leftovers:
            source = entry["md_path"]
            dest = folder_path / source.name
            try:
                if dest.exists():
                    dest.unlink()
                shutil.move(str(source), str(dest))
            except Exception:
                logger.exception("Failed to move %s to %s", source, dest)
                continue
            entry["md_path"] = dest
            assigned.add(entry["number"])
            moved_units.append(
                {
                    "number": entry["number"],
                    "summary": unit_summaries.get(entry["number"], ""),
                    "relative_path": dest.relative_to(output_dir).as_posix(),
                }
            )
        if moved_units:
            category_details.append(
                {
                    "name": "Uncategorized",
                    "description": "自動分類でカテゴリが決まらなかったページです。",
                    "folder": folder_name,
                    "pages": moved_units,
                }
            )

    sections: list[str] = []
    for cat in category_details:
        description = cat["description"].strip()
        header = f"## {cat['name']}\n保存フォルダ: `{cat['folder']}/`\n"
        if description:
            header += f"{description}\n"
        lines = [
            f"- {prompts['unit_label_en']} {info['number']:03}: {info['summary'] or '(要約なし)'} (ファイル: `{info['relative_path']}`)"
            for info in cat["pages"]
        ]
        sections.append(header + "\n".join(lines))

    if not sections:
        return f"{prompts['unit_label_jp']}分類は生成されませんでした。"

    return "\n\n".join(sections)


def _process_document_images(
    *,
    model: BaseChatModel,
    output_dir: Path,
    document_path: Path,
    image_results: list[tuple[int, Path, str]],
    prompts: DocumentProcessingPrompts,
    extra_text_lookup: dict[int, str] | None = None,
    delete_after: list[Path] | None = None,
    link_base_url: str | None = None,
) -> None:
    if not image_results:
        logger.info("No images provided for document %s; skipping extraction", document_path.name)
        _write_summary(output_dir, document_path.name, f"{prompts['unit_label_jp']}画像を生成できませんでした。")
        return

    lookup = extra_text_lookup or {}
    messages = [
        _build_messages_for_unit(number, data_url, prompts, extra_text=lookup.get(number))
        for number, _, data_url in image_results
    ]

    try:
        responses = model.batch(messages, config={"max_concurrency": 4})
    except Exception as exc:
        logger.exception("LangChain batch inference failed for %s: %s", document_path.name, exc)
        _write_summary(output_dir, document_path.name, f"抽出処理中にエラーが発生しました。\n\n詳細: {exc}")
        return

    unit_entries: list[dict] = []
    unit_prefix = prompts["unit_prefix"]
    unit_label = prompts["unit_label_en"]
    for (unit_number, image_path, _), response in zip(image_results, responses):
        unit_text = _message_content_to_text(response).strip()
        if not unit_text:
            unit_text = "_(No readable content detected for this image)_"
        unit_md_path = output_dir / f"{unit_prefix}_{unit_number:03}.md"
        try:
            header = f"# {unit_label} {unit_number}\n\n"
            link_line = f"ページリンク：{link_base_url}#page={unit_number}\n\n" if link_base_url else ""
            body = f"{unit_text}\n"
            unit_md_path.write_text(header + link_line + body, encoding="utf-8")
        except Exception as write_exc:
            logger.error(
                "Failed to write markdown for %s %s of %s: %s",
                unit_label,
                unit_number,
                document_path.name,
                write_exc,
            )
            _write_summary(
                output_dir,
                document_path.name,
                f"{prompts['unit_label_jp']} {unit_number} のMarkdownを書き込めませんでした。\n\n詳細: {write_exc}",
            )
            return
        unit_entries.append(
            {
                "number": unit_number,
                "md_path": unit_md_path,
                "image_path": image_path,
                "text": unit_text,
            }
        )

    summary_warning: str | None = None
    try:
        unit_summaries = _summarize_units(model, unit_entries, prompts)
    except Exception as exc:
        logger.exception("Failed to summarize units for %s: %s", document_path.name, exc)
        summary_warning = f"{prompts['unit_label_jp']}要約の生成に失敗しました: {exc}"
        unit_summaries = {entry["number"]: "(要約生成に失敗しました)" for entry in unit_entries}

    try:
        categories = _categorize_units(model, unit_summaries, prompts)
    except Exception as exc:
        logger.exception("Failed to categorize units for %s: %s", document_path.name, exc)
        summary_warning = f"{prompts['unit_label_jp']}分類の生成に失敗しました: {exc}"
        categories = [{"name": "Uncategorized", "description": "", "pages": sorted(unit_summaries.keys())}]

    organization_summary = _organize_units_by_category(
        output_dir, unit_entries, categories, unit_summaries, prompts
    )

    for entry in unit_entries:
        image_path = entry["image_path"]
        try:
            if image_path.exists():
                image_path.unlink()
        except Exception:
            logger.exception("Failed to delete image %s", image_path)

    intro_lines = [
        f"{prompts['unit_label_jp']}画像（PNG）はカテゴリ整理後に削除済みです。",
        "Markdownは自動生成されたカテゴリごとのサブフォルダに移動済みです。",
    ]
    if summary_warning:
        intro_lines.append(f"注意: {summary_warning}")

    summary_message = "\n".join(intro_lines) + "\n\n" + organization_summary
    _write_summary(output_dir, document_path.name, summary_message)
    summary_md_path = output_dir / "summary.md"
    try:
        if summary_md_path.exists():
            summary_md_path.unlink()
    except Exception:
        logger.exception("Failed to delete summary markdown %s", summary_md_path)

    # Delete any additional paths (e.g., original PPTX) now that processing is complete
    if delete_after:
        for target in delete_after:
            try:
                if target.exists():
                    target.unlink()
            except Exception:
                logger.exception("Failed to delete original file %s", target)


def _build_langchain_model(profile: LLMProfile) -> BaseChatModel:
    errors: list[str] = []

    def _openai_model() -> BaseChatModel:
        cfg = profile.openai
        if not cfg.apiKey:
            raise RuntimeError("OpenAI API key is not configured")
        model_name = cfg.model or os.environ.get("OPENAI_MODEL") or "gpt-4o-mini"
        return ChatOpenAI(
            api_key=cfg.apiKey,
            base_url=(cfg.baseUrl or None),
            model=model_name,
            temperature=0.0,
            max_retries=2,
        )

    def _azure_model() -> BaseChatModel:
        cfg = profile.azure
        if not (cfg.apiKey and cfg.endpoint and cfg.deployment):
            raise RuntimeError("Azure OpenAI is not fully configured")
        return AzureChatOpenAI(
            api_key=cfg.apiKey,
            azure_endpoint=cfg.endpoint,
            azure_deployment=cfg.deployment,
            api_version=(cfg.apiVersion or "2024-02-15-preview"),
            temperature=0.0,
            max_retries=2,
        )

    builders = (_azure_model, _openai_model) if profile.preferred == "azure" else (_openai_model, _azure_model)
    for builder in builders:
        try:
            return builder()
        except Exception as exc:
            errors.append(str(exc))
            continue
    raise RuntimeError("; ".join(errors) or "Failed to initialize any LLM client")


def _process_pdf_upload(
    username: str,
    pdf_path_str: str,
    source_url: str | None = None,
    delete_after: list[str] | None = None,
    progress_rel_key: str | None = None,
) -> None:
    pdf_path = Path(pdf_path_str)
    if not pdf_path.exists() or not pdf_path.is_file():
        logger.warning("PDF path %s does not exist; skipping extraction", pdf_path_str)
        return

    output_dir = _ensure_output_dir(pdf_path)
    try:
        # Initialize progress for multi-step pipeline
        default_rel = pdf_path.relative_to(BASE_DIR / username).as_posix()
    except Exception:
        default_rel = pdf_path.name
    rel = progress_rel_key or default_rel
    _set_progress(username, rel, stage="preprocess", percent=0.05, message="PDFを準備中", done=False)
    _cleanup_previous_outputs(output_dir, preserve={pdf_path})

    profile = _load_llm_profile(username)
    try:
        model = _build_langchain_model(profile)
    except Exception as exc:
        logger.warning("Skipping PDF extraction for %s: %s", pdf_path.name, exc)
        _write_summary(output_dir, pdf_path.name, f"PDF解析に必要なLLM設定が見つかりませんでした。\n\n詳細: {exc}")
        return

    page_results: list[tuple[int, Path, str]] = []
    try:
        with pdfplumber.open(pdf_path) as pdf:
            if not pdf.pages:
                logger.info("PDF %s has no pages; skipping extraction", pdf_path.name)
                _write_summary(output_dir, pdf_path.name, "PDFにページが含まれていませんでした。")
                _set_progress(username, rel, stage="done", percent=1.0, message="完了", done=True)
                _clear_progress(username, rel)
                return
            total_pages = len(pdf.pages)
            for idx, page in enumerate(pdf.pages, start=1):
                page_image = page.to_image(resolution=300)
                pil_image_raw = getattr(page_image, "original", None)
                if pil_image_raw is None or not isinstance(pil_image_raw, Image.Image):
                    raise RuntimeError("pdfplumber failed to provide a PIL image for page conversion")
                pil_image = pil_image_raw.convert("RGB")
                image_path = output_dir / f"page_{idx:03}.png"
                pil_image.save(str(image_path), format="PNG")
                data_url = _pil_image_to_data_url(pil_image)
                page_results.append((idx, image_path, data_url))
                # Update progress during rasterization (10% -> 40%)
                frac = 0.1 + 0.3 * (idx / max(1, total_pages))
                _set_progress(username, rel, stage="rasterize", percent=frac, message=f"ページ画像化 {idx}/{total_pages}")
    except Exception as exc:
        logger.exception("Failed to rasterize PDF %s: %s", pdf_path.name, exc)
        _write_summary(output_dir, pdf_path.name, f"PDFのページ画像化に失敗しました。\n\n詳細: {exc}")
        _set_progress(username, rel, stage="error", percent=1.0, message="ページ画像化に失敗", done=True)
        _clear_progress(username, rel)
        return

    if not page_results:
        logger.info("No page images generated for %s; skipping extraction", pdf_path.name)
        _write_summary(output_dir, pdf_path.name, "ページ画像を生成できませんでした。")
        return

    pdf_prompts: DocumentProcessingPrompts = {
        "doc_label_jp": "PDF文書",
        "unit_label_en": "Page",
        "unit_label_jp": "ページ",
        "unit_prefix": "page",
        "extraction_system_prompt": PDF_EXTRACTION_SYSTEM_PROMPT,
        "extraction_instruction": PDF_EXTRACTION_PAGE_INSTRUCTIONS,
        "summary_system_prompt": PAGE_SUMMARY_SYSTEM_PROMPT,
        "category_system_prompt": CATEGORY_SYSTEM_PROMPT,
    }

    _set_progress(username, rel, stage="extract", percent=0.45, message="Markdown抽出中")
    _process_document_images(
        model=model,
        output_dir=output_dir,
        document_path=pdf_path,
        image_results=page_results,
        prompts=pdf_prompts,
        delete_after=[Path(p) for p in delete_after] if delete_after else None,
        link_base_url=source_url,
    )
    _set_progress(username, rel, stage="categorize", percent=0.85, message="カテゴリ分類と移動中")
    # _process_document_images performs categorization and moves; after it returns, consider done
    _set_progress(username, rel, stage="done", percent=1.0, message="完了", done=True)
    _clear_progress(username, rel)


def _process_pptx_upload(username: str, pptx_path_str: str) -> None:
    pptx_path = Path(pptx_path_str)
    if not pptx_path.exists() or not pptx_path.is_file():
        logger.warning("PPTX path %s does not exist; skipping extraction", pptx_path_str)
        return

    output_dir = _ensure_output_dir(pptx_path)
    try:
        rel = pptx_path.relative_to(BASE_DIR / username).as_posix()
    except Exception:
        rel = pptx_path.name
    _set_progress(username, rel, stage="preprocess", percent=0.05, message="PPTXを準備中", done=False)
    _cleanup_previous_outputs(output_dir, preserve={pptx_path})

    if PptxPresentation is None or MSO_SHAPE_TYPE is None:
        message = "PPTX解析に必要なライブラリ (python-pptx) がインストールされていません。"
        logger.warning(message)
        _write_summary(output_dir, pptx_path.name, message)
        return

    # First, try high-fidelity conversion via LibreOffice/soffice/unoconv
    pdf_converted_path = _convert_pptx_to_pdf(pptx_path)
    if pdf_converted_path is not None and pdf_converted_path.exists():
        _set_progress(username, rel, stage="convert", percent=0.2, message="PDFへ変換中")
        # Analyze the generated PDF and delete the original PPTX after categorization
        _process_pdf_upload(username, str(pdf_converted_path), source_url=None, delete_after=[str(pptx_path)])
        _set_progress(username, rel, stage="done", percent=1.0, message="完了", done=True)
        _clear_progress(username, rel)
        return

    # Fallback: proceed with python-pptx based slide rasterization
    profile = _load_llm_profile(username)
    try:
        model = _build_langchain_model(profile)
    except Exception as exc:
        logger.warning("Skipping PPTX extraction for %s: %s", pptx_path.name, exc)
        _write_summary(
            output_dir,
            pptx_path.name,
            f"PPTX解析に必要なLLM設定が見つかりませんでした。\n\n詳細: {exc}",
        )
        return

    try:
        presentation = PptxPresentation(str(pptx_path))
    except Exception as exc:
        logger.exception("Failed to load PPTX %s: %s", pptx_path.name, exc)
        _write_summary(output_dir, pptx_path.name, f"PPTXの読み込みに失敗しました。\n\n詳細: {exc}")
        _set_progress(username, rel, stage="error", percent=1.0, message="PPTX読み込みに失敗", done=True)
        _clear_progress(username, rel)
        return

    slide_results: list[tuple[int, Path, str]] = []
    extra_text_lookup: dict[int, str] = {}
    try:
        slides_iter = list(presentation.slides)
    except Exception as exc:
        logger.exception("Failed to enumerate slides for %s: %s", pptx_path.name, exc)
        _write_summary(output_dir, pptx_path.name, f"PPTXのスライド取得に失敗しました。\n\n詳細: {exc}")
        _set_progress(username, rel, stage="error", percent=1.0, message="スライド取得に失敗", done=True)
        _clear_progress(username, rel)
        return

    if not slides_iter:
        logger.info("PPTX %s has no slides; skipping extraction", pptx_path.name)
        _write_summary(output_dir, pptx_path.name, "PPTXにスライドが含まれていませんでした。")
        return

    for idx, slide in enumerate(slides_iter, start=1):
        try:
            image = _render_pptx_slide(slide, presentation)
        except Exception as exc:
            logger.exception("Failed to render slide %s of %s: %s", idx, pptx_path.name, exc)
            _write_summary(output_dir, pptx_path.name, f"スライド {idx} の画像生成に失敗しました。\n\n詳細: {exc}")
            _set_progress(username, rel, stage="error", percent=1.0, message="スライド画像生成に失敗", done=True)
            _clear_progress(username, rel)
            return
        image_path = output_dir / f"slide_{idx:03}.png"
        try:
            image.save(str(image_path), format="PNG")
        except Exception as exc:
            logger.exception("Failed to save PNG for slide %s of %s: %s", idx, pptx_path.name, exc)
            _write_summary(output_dir, pptx_path.name, f"スライド {idx} のPNG保存に失敗しました。\n\n詳細: {exc}")
            _set_progress(username, rel, stage="error", percent=1.0, message="PNG保存に失敗", done=True)
            _clear_progress(username, rel)
            return
        data_url = _pil_image_to_data_url(image)
        slide_results.append((idx, image_path, data_url))
        _set_progress(username, rel, stage="rasterize", percent=0.2 + 0.4 * (idx / max(1, len(slides_iter))), message=f"スライド画像化 {idx}/{len(slides_iter)}")

        try:
            context_text = _extract_pptx_slide_context(slide)
            if context_text:
                extra_text_lookup[idx] = context_text
        except Exception:
            logger.exception("Failed to extract context text for slide %s of %s", idx, pptx_path.name)

    pptx_prompts: DocumentProcessingPrompts = {
        "doc_label_jp": "PPTX資料",
        "unit_label_en": "Slide",
        "unit_label_jp": "スライド",
        "unit_prefix": "slide",
        "extraction_system_prompt": PPTX_EXTRACTION_SYSTEM_PROMPT,
        "extraction_instruction": PPTX_EXTRACTION_SLIDE_INSTRUCTIONS,
        "summary_system_prompt": _build_summary_system_prompt("PPTX資料", "スライド"),
        "category_system_prompt": _build_category_system_prompt("PPTX資料", "スライド"),
    }

    _set_progress(username, rel, stage="extract", percent=0.65, message="Markdown抽出中")
    _process_document_images(
        model=model,
        output_dir=output_dir,
        document_path=pptx_path,
        image_results=slide_results,
        prompts=pptx_prompts,
        extra_text_lookup=extra_text_lookup if extra_text_lookup else None,
        delete_after=[pptx_path],
    )
    _set_progress(username, rel, stage="categorize", percent=0.9, message="カテゴリ分類と移動中")
    _set_progress(username, rel, stage="done", percent=1.0, message="完了", done=True)
    _clear_progress(username, rel)

    # Try to assemble a PDF from the slide PNGs so that a PDF remains even in fallback
    try:
        png_paths = [path for _, path, _ in slide_results]
        if png_paths:
            fallback_pdf = output_dir / f"{pptx_path.stem}.pdf"
            _combine_images_to_pdf(png_paths, fallback_pdf)
    except Exception:
        logger.exception("Failed to build fallback PDF for %s", pptx_path.name)

def _convert_pptx_to_pdf(pptx_path: Path) -> Path | None:
    """Convert PPTX to PDF using available system tools. Returns PDF path on success."""
    out_dir = pptx_path.parent
    expected_pdf = out_dir / f"{pptx_path.stem}.pdf"

    def run_cmd(cmd: list[str]) -> bool:
        try:
            result = subprocess.run(
                cmd,
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
                check=True,
            )
            return result.returncode == 0
        except Exception:
            logger.exception("PPTX->PDF conversion command failed: %s", " ".join(cmd))
            return False

    # Windows: prefer native PowerPoint COM automation if available
    try:
        if platform.system().lower().startswith("win"):
            # Try win32com first
            try:
                import win32com.client  # type: ignore
                pp = win32com.client.Dispatch("PowerPoint.Application")
                pp.Visible = 0
                try:
                    presentation = pp.Presentations.Open(str(pptx_path), WithWindow=False)
                except Exception:
                    presentation = pp.Presentations.Open(str(pptx_path))
                try:
                    # 32 = ppSaveAsPDF
                    presentation.SaveAs(str(expected_pdf), 32)
                finally:
                    try:
                        presentation.Close()
                    except Exception:
                        pass
                    try:
                        pp.Quit()
                    except Exception:
                        pass
                if expected_pdf.exists():
                    return expected_pdf
            except Exception:
                # Try comtypes as a fallback on Windows
                try:
                    import comtypes.client  # type: ignore
                    pp = comtypes.client.CreateObject("PowerPoint.Application")
                    pp.Visible = 0
                    presentation = None
                    try:
                        # Some environments support named arg
                        presentation = pp.Presentations.Open(str(pptx_path), WithWindow=False)
                    except Exception:
                        try:
                            # Positional fallback: (FileName, ReadOnly, Untitled, WithWindow)
                            presentation = pp.Presentations.Open(str(pptx_path), False, False, False)
                        except Exception:
                            pass
                    if presentation is not None:
                        try:
                            presentation.SaveAs(str(expected_pdf), 32)
                        finally:
                            try:
                                presentation.Close()
                            except Exception:
                                pass
                            try:
                                pp.Quit()
                            except Exception:
                                pass
                        if expected_pdf.exists():
                            return expected_pdf
                except Exception:
                    logger.exception("PowerPoint COM automation failed for %s", pptx_path.name)
    except Exception:
        logger.exception("Windows PowerPoint-based conversion failed for %s", pptx_path.name)

    # Prefer soffice if available
    soffice = shutil.which("soffice")
    # On Windows, also try common install paths for LibreOffice
    if not soffice and platform.system().lower().startswith("win"):
        candidates = [
            Path("C:/Program Files/LibreOffice/program/soffice.exe"),
            Path("C:/Program Files (x86)/LibreOffice/program/soffice.exe"),
        ]
        for c in candidates:
            if c.exists():
                soffice = str(c)
                break
    if soffice:
        # Convert into the same directory
        ok = run_cmd([soffice, "--headless", "--nologo", "--convert-to", "pdf", "--outdir", str(out_dir), str(pptx_path)])
        if ok and expected_pdf.exists():
            return expected_pdf

    libreoffice = shutil.which("libreoffice")
    if libreoffice:
        ok = run_cmd([libreoffice, "--headless", "--nologo", "--convert-to", "pdf", "--outdir", str(out_dir), str(pptx_path)])
        if ok and expected_pdf.exists():
            return expected_pdf

    unoconv = shutil.which("unoconv")
    if unoconv:
        ok = run_cmd([unoconv, "-f", "pdf", "-o", str(expected_pdf), str(pptx_path)])
        if ok and expected_pdf.exists():
            return expected_pdf

    logger.warning("No suitable PPTX->PDF converter found or conversion failed for %s", pptx_path.name)
    return None

def _combine_images_to_pdf(image_paths: list[Path], pdf_path: Path) -> None:
    images: list[Image.Image] = []
    for p in image_paths:
        try:
            img = Image.open(str(p))
            if img.mode != "RGB":
                img = img.convert("RGB")
            images.append(img)
        except Exception:
            logger.exception("Failed to open image %s for PDF assembly", p)
    if not images:
        return
    first, rest = images[0], images[1:]
    try:
        first.save(str(pdf_path), save_all=True, append_images=rest, format="PDF")
    except Exception:
        logger.exception("Failed to write assembled PDF %s", pdf_path)
# External service endpoints
MCP_SERVER_BASE = os.environ.get("MCP_SERVER_BASE", "http://localhost:3030")

# -------- LLM per-user configuration --------

class OpenAIConfig(BaseModel):
    apiKey: str = ""
    baseUrl: str | None = None  # optional override; default SDK base when None/empty
    model: str | None = None  # optional; defaults to "gpt-4o-mini" when not set

class AzureConfig(BaseModel):
    apiKey: str = ""
    endpoint: str | None = None  # e.g. https://<resource>.openai.azure.com
    deployment: str | None = None  # Azure deployment name
    apiVersion: str | None = None  # e.g. 2024-02-15-preview

class LLMProfile(BaseModel):
    preferred: Literal["openai", "azure"] = "openai"
    openai: OpenAIConfig = OpenAIConfig()
    azure: AzureConfig = AzureConfig()

def _get_user_record(users: dict, username_or_canon: str) -> tuple[str, dict]:
    canon = _canonicalize_username(username_or_canon)
    key = canon if canon in users else next((k for k in users.keys() if k.lower() == canon), canon)
    return key, users.get(key) or {}

def _load_llm_profile(username: str) -> LLMProfile:
    users = _load_users()
    _, record = _get_user_record(users, username)
    cfg = (record.get("llm") or {})
    try:
        return LLMProfile(**cfg)  # type: ignore[arg-type]
    except Exception:
        return LLMProfile()

def _save_llm_profile(username: str, profile: LLMProfile) -> None:
    users = _load_users()
    key, record = _get_user_record(users, username)
    record = dict(record)
    record["llm"] = profile.model_dump()
    users[key] = record
    _save_users(users)

class LLMProfileResponse(LLMProfile):
    pass

@app.get("/api/profile/llm", response_model=LLMProfileResponse)
def get_llm_profile(request: Request):
    username, _ = require_user(request)
    return _load_llm_profile(username)

@app.post("/api/profile/llm", response_model=LLMProfileResponse)
def update_llm_profile(request: Request, body: LLMProfile):
    username, _ = require_user(request)
    # Normalize empty strings to None for URLs/optional fields
    normalized = LLMProfile(
        preferred=body.preferred,
        openai=OpenAIConfig(
            apiKey=(body.openai.apiKey or "").strip(),
            baseUrl=(body.openai.baseUrl or None) or None,
            model=(body.openai.model or None) or None,
        ),
        azure=AzureConfig(
            apiKey=(body.azure.apiKey or "").strip(),
            endpoint=(body.azure.endpoint or None) or None,
            deployment=(body.azure.deployment or None) or None,
            apiVersion=(body.azure.apiVersion or None) or None,
        ),
    )
    _save_llm_profile(username, normalized)
    return normalized

# Trailing-slash aliases for compatibility
@app.get("/api/profile/llm/", response_model=LLMProfileResponse)
def get_llm_profile_slash(request: Request):
    return get_llm_profile(request)

@app.post("/api/profile/llm/", response_model=LLMProfileResponse)
def update_llm_profile_slash(request: Request, body: LLMProfile):
    return update_llm_profile(request, body)

# Secret for signing share tokens (fall back to a constant for local dev)
SHARE_SECRET = (os.environ.get("SHARE_SECRET") or "dev-secret").encode("utf-8")

def _b64url(data: bytes) -> str:
    return base64.urlsafe_b64encode(data).rstrip(b"=").decode("ascii")

def _b64url_decode(data: str) -> bytes:
    pad = '=' * (-len(data) % 4)
    return base64.urlsafe_b64decode(data + pad)

def sign_share_token(path: str, expires_epoch: int) -> str:
    msg = f"{path}\n{expires_epoch}".encode("utf-8")
    sig = hmac.new(SHARE_SECRET, msg, hashlib.sha256).digest()
    return _b64url(sig)

def verify_share_token(path: str, expires_epoch: int, token: str) -> bool:
    try:
        sig = _b64url_decode(token)
    except Exception:
        return False
    msg = f"{path}\n{expires_epoch}".encode("utf-8")
    expected = hmac.new(SHARE_SECRET, msg, hashlib.sha256).digest()
    return hmac.compare_digest(sig, expected)


# CORS for Vite dev server
app.add_middleware(
    CORSMiddleware,
    allow_origins=[
        "http://localhost:5173",
        "http://127.0.0.1:5173",
        "null",  # allow file:// origin for local testing
    ],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)


# Disable caching of API responses to avoid storing sensitive authenticated data
@app.middleware("http")
async def add_no_store_headers(request: Request, call_next):
    response = await call_next(request)
    try:
        response.headers["Cache-Control"] = "no-store, no-cache, must-revalidate, max-age=0"
        response.headers["Pragma"] = "no-cache"
        response.headers["Expires"] = "0"
    except Exception:
        pass
    return response


def safe_join(base: Path, *paths: str) -> Path:
    candidate = (base.joinpath(*paths)).resolve()
    try:
        candidate.relative_to(base)
    except ValueError:
        raise HTTPException(status_code=400, detail="Invalid path")
    return candidate


def safe_join_leaf(base: Path, rel_path: str) -> Path:
    """Join a relative path under base but do not resolve the final leaf component.

    This preserves symlink behavior for the final path while still ensuring the
    parent directory is safely contained within the base directory.
    """
    if rel_path.strip() == "":
        raise HTTPException(status_code=400, detail="Invalid path")
    # Normalize separators and trim surrounding slashes
    normalized = rel_path.replace("\\", "/").strip("/")
    # Compute parent under base and resolve to validate containment
    parent_rel = Path(normalized).parent
    leaf = Path(normalized).name
    validate_name(leaf)
    parent_abs = (base / parent_rel).resolve()
    try:
        parent_abs.relative_to(base)
    except ValueError:
        raise HTTPException(status_code=400, detail="Invalid path")
    return parent_abs / leaf


def stat_entry(path: Path, rel_base: Path) -> Entry:
    is_link = path.is_symlink()
    try:
        st = path.stat()
        is_dir = path.is_dir()
        entry_type: Literal["file", "directory"] = "directory" if is_dir else "file"
        size_val = None if is_dir else st.st_size
    except FileNotFoundError:
        # Likely a broken symlink or concurrently removed entry. Use lstat for metadata
        try:
            st = path.lstat()  # type: ignore[attr-defined]
        except Exception:
            # If even lstat fails (deleted), propagate to caller to optionally skip
            raise
        entry_type = "file"  # fall back to file to avoid navigating into broken links
        size_val = None
    rel_path = path.relative_to(rel_base)
    return {
        "name": path.name,
        "path": str(rel_path).replace(os.sep, "/"),
        "type": entry_type,
        "size": size_val,
        "mtime": st.st_mtime,
        "isLink": is_link,
    }


_INVALID_NAME_PATTERN = re.compile(r"[\\/\0]")


def validate_name(name: str) -> None:
    if not name or name.strip() == "":
        raise HTTPException(status_code=400, detail="Name is required")
    if name in {".", ".."}:
        raise HTTPException(status_code=400, detail="Invalid name")
    if _INVALID_NAME_PATTERN.search(name):
        raise HTTPException(status_code=400, detail="Name contains invalid characters")


# -------------------- Simple file-based user auth --------------------
USERS_DIR = (BASE_DIR / "_users")
USERS_FILE = (USERS_DIR / "users.json")

def _ensure_users_file() -> None:
    try:
        USERS_DIR.mkdir(parents=True, exist_ok=True)
        if not USERS_FILE.exists():
            USERS_FILE.write_text("{}", encoding="utf-8")
    except Exception:
        raise HTTPException(status_code=500, detail="Failed to initialize users store")

def _load_users() -> dict:
    _ensure_users_file()
    try:
        return json.loads(USERS_FILE.read_text(encoding="utf-8"))
    except Exception:
        return {}

def _save_users(users: dict) -> None:
    _ensure_users_file()
    try:
        USERS_FILE.write_text(json.dumps(users, ensure_ascii=False, indent=2), encoding="utf-8")
    except Exception:
        raise HTTPException(status_code=500, detail="Failed to save users store")

def _hash_password(username: str, password: str) -> str:
    # Simple SHA256-based hash with username-specific salt and server-side pepper
    data = f"{username}\n{password}".encode("utf-8") + SHARE_SECRET
    return hashlib.sha256(data).hexdigest()

def _token_sign(payload: dict) -> str:
    body = json.dumps(payload, separators=(",", ":"), ensure_ascii=False).encode("utf-8")
    b64 = _b64url(body)
    sig = hmac.new(SHARE_SECRET, b64.encode("ascii"), hashlib.sha256).digest()
    return f"{b64}.{_b64url(sig)}"

def _token_verify(token: str) -> dict:
    try:
        b64, sig = token.split(".")
    except ValueError:
        raise HTTPException(status_code=401, detail="Invalid token")
    expected = _b64url(hmac.new(SHARE_SECRET, b64.encode("ascii"), hashlib.sha256).digest())
    if not hmac.compare_digest(sig, expected):
        raise HTTPException(status_code=401, detail="Invalid token")
    payload = json.loads(_b64url_decode(b64).decode("utf-8", errors="replace"))
    # Optional expiry
    exp = payload.get("exp")
    if isinstance(exp, int) and int(time.time()) > exp:
        raise HTTPException(status_code=401, detail="Token expired")
    return payload

def _get_auth_header(request: Request) -> str:
    auth = request.headers.get("authorization") or request.headers.get("Authorization") or ""
    return auth


def _canonicalize_username(username: str) -> str:
    """Return a canonical username representation used for storage and directory naming.

    We normalize to lowercase to avoid case-variant duplication across browsers/OSes.
    """
    return (username or "").strip().lower()


def _resolve_user_base_dir(preferred_username: str) -> Path:
    """Resolve a user's base directory with case-insensitive fallback.

    - Prefer the canonical directory (lowercased)
    - If it doesn't exist, but a case-variant exists (e.g. "Sho"), use that instead
    - Otherwise return the canonical path (caller may create it)
    """
    canonical = _canonicalize_username(preferred_username)
    canonical_path = (BASE_DIR / canonical)
    if canonical_path.exists() and canonical_path.is_dir():
        return canonical_path

    # Fallback: search for an existing case-variant under BASE_DIR
    try:
        for child in BASE_DIR.iterdir():
            try:
                if child.is_dir() and child.name.lower() == canonical and child.name not in {".", "..", "_users"}:
                    return child
            except Exception:
                continue
    except Exception:
        # If listing fails, just return canonical
        pass

    return canonical_path

def require_user(request: Request) -> tuple[str, Path]:
    auth = _get_auth_header(request)
    if not auth.startswith("Bearer "):
        raise HTTPException(status_code=401, detail="Unauthorized")
    token = auth[7:].strip()
    payload = _token_verify(token)
    raw_username = str(payload.get("sub") or "").strip()
    if not raw_username:
        raise HTTPException(status_code=401, detail="Unauthorized")
    canonical_username = _canonicalize_username(raw_username)
    # Enforce token version (revocation) against persisted users store
    try:
        users = _load_users()
        key = canonical_username if canonical_username in users else next((k for k in users.keys() if k.lower() == canonical_username), None)
        if not key:
            raise HTTPException(status_code=401, detail="Unauthorized")
        stored_version = int((users.get(key) or {}).get("tokenVersion", 1))
        token_version = int((payload.get("v") if isinstance(payload.get("v"), int) else 0) or 0) or 1
        if token_version != stored_version:
            raise HTTPException(status_code=401, detail="Unauthorized")
    except HTTPException:
        raise
    except Exception:
        # Deny on any unexpected error
        raise HTTPException(status_code=401, detail="Unauthorized")
    # Resolve to an existing case-variant folder when present to avoid user data appearing missing
    user_base = _resolve_user_base_dir(canonical_username)
    user_base.mkdir(parents=True, exist_ok=True)
    return canonical_username, user_base

class AuthRegisterRequest(BaseModel):
    username: str
    password: str

class AuthLoginRequest(BaseModel):
    username: str
    password: str

class AuthMeResponse(BaseModel):
    username: str

@app.post("/api/auth/register", response_model=AuthMeResponse)
def auth_register(req: AuthRegisterRequest):
    username = req.username.strip()
    if not username or "/" in username or "\\" in username or username in {".", "..", "_users"}:
        raise HTTPException(status_code=400, detail="Invalid username")
    if not req.password:
        raise HTTPException(status_code=400, detail="Password required")
    users = _load_users()
    canon = _canonicalize_username(username)
    # Prevent duplicates across case variants
    if canon in users or any(k.lower() == canon for k in users.keys()):
        raise HTTPException(status_code=409, detail="User exists")
    users[canon] = {"password": _hash_password(canon, req.password), "created": int(time.time()), "tokenVersion": 1}
    _save_users(users)
    # Create user's folder
    user_dir = BASE_DIR / canon
    try:
        user_dir.mkdir(parents=True, exist_ok=True)
    except Exception:
        raise HTTPException(status_code=500, detail="Failed to create user directory")
    # Return minimal profile (frontend will store token from header)
    return AuthMeResponse(username=canon)

@app.post("/api/auth/login", response_model=AuthMeResponse)
def auth_login(req: AuthLoginRequest):
    users = _load_users()
    provided = req.username.strip()
    canon = _canonicalize_username(provided)
    # Prefer canonical entry; fallback to any case-variant key
    found_key = canon if canon in users else next((k for k in users.keys() if k.lower() == canon), None)
    if not found_key:
        raise HTTPException(status_code=401, detail="Invalid credentials")
    user = users.get(found_key) or {}
    if user.get("password") != _hash_password(found_key, req.password):
        raise HTTPException(status_code=401, detail="Invalid credentials")
    # Always return canonical username for consistency across clients
    return AuthMeResponse(username=_canonicalize_username(found_key))

@app.get("/api/auth/token")
def auth_token(username: str = Query(...)):
    # Utility endpoint to mint a token after successful login/register
    # Client calls this immediately to obtain a signed token
    users = _load_users()
    canon = _canonicalize_username(username)
    # Accept legacy case keys
    if not (canon in users or any(k.lower() == canon for k in users.keys())):
        raise HTTPException(status_code=404, detail="User not found")
    # Include current tokenVersion so server can revoke old tokens
    key = canon if canon in users else next((k for k in users.keys() if k.lower() == canon), canon)
    version = int((users.get(key) or {}).get("tokenVersion", 1))
    # Issue a stable token (no expiry) so it remains constant until rotation
    payload = {"sub": canon, "v": version}
    return {"token": _token_sign(payload)}

@app.post("/api/auth/token/rotate")
def auth_token_rotate(request: Request):
    # Rotate user's token by incrementing tokenVersion; returns a new signed token
    username, _ = require_user(request)
    users = _load_users()
    key = username if username in users else next((k for k in users.keys() if k.lower() == username), username)
    record = users.get(key) or {}
    current = int(record.get("tokenVersion", 1))
    record["tokenVersion"] = current + 1
    users[key] = record
    _save_users(users)
    # New token (no expiry) tied to incremented version
    payload = {"sub": _canonicalize_username(key), "v": int(record["tokenVersion"])}
    return {"token": _token_sign(payload)}

@app.get("/api/auth/me", response_model=AuthMeResponse)
def auth_me(request: Request):
    username, _ = require_user(request)
    return AuthMeResponse(username=username)


@app.get("/api/browse", response_model=List[Entry])
def browse(request: Request, path: str = Query("", description="Relative path under user's base directory")):
    """List files and folders under the provided path (relative to BASE_DIR)."""
    _, user_base = require_user(request)
    target = safe_join(user_base, path)
    if not target.exists():
        raise HTTPException(status_code=404, detail="Path not found")
    if not target.is_dir():
        raise HTTPException(status_code=400, detail="Path is not a directory")

    entries: List[Entry] = []
    try:
        for child in sorted(target.iterdir(), key=lambda p: (p.is_file(), p.name.lower())):
            try:
                entries.append(stat_entry(child, user_base))
            except FileNotFoundError:
                # Skip entries that disappeared or broken links we couldn't lstat
                continue
    except PermissionError:
        raise HTTPException(status_code=403, detail="Permission denied")

    return entries


@app.get("/api/download")
def download(request: Request, path: str = Query(..., description="Relative file path under user's base directory"), token: str | None = Query(None), inline: bool = Query(False, description="Set Content-Disposition to inline for browser preview")):
    # Allow header auth or token query param
    username = None
    user_base: Path | None = None
    if token:
        payload = _token_verify(token)
        username = str(payload.get("sub") or "").strip()
        if not username:
            raise HTTPException(status_code=401, detail="Unauthorized")
        user_base = (BASE_DIR / username)
        user_base.mkdir(parents=True, exist_ok=True)
    else:
        _, user_base = require_user(request)
    file_path = safe_join(user_base, path)
    if not file_path.exists() or not file_path.is_file():
        raise HTTPException(status_code=404, detail="File not found")

    media_type, _ = mimetypes.guess_type(str(file_path))
    response = FileResponse(
        str(file_path),
        media_type=media_type or "application/octet-stream",
    )
    # Build a safe Content-Disposition header
    disp_type = "inline" if inline else "attachment"
    filename = file_path.name
    try:
        ascii_name = filename.encode("latin-1", "strict").decode("latin-1")
        cd = f"{disp_type}; filename=\"{ascii_name}\""
    except UnicodeEncodeError:
        # Fallback for non-ASCII filenames per RFC 5987 (UTF-8 percent-encoded)
        encoded = urllib.parse.quote(filename, safe="")
        cd = f"{disp_type}; filename*=UTF-8''{encoded}"
    response.headers["Content-Disposition"] = cd
    return response


class ReadFileResponse(BaseModel):
    path: str
    content: str
    encoding: str


@app.get("/api/read", response_model=ReadFileResponse)
def read_file_contents(
    request: Request,
    path: str = Query(..., description="Relative file path under user's base directory"),
    encoding: str = Query("utf-8", description="Text encoding to use when reading the file"),
):
    """Read a text file under BASE_DIR and return its contents as JSON.

    The path is constrained to the configured base directory via safe_join.
    """
    _, user_base = require_user(request)
    file_path = safe_join(user_base, path)
    if not file_path.exists() or not file_path.is_file():
        raise HTTPException(status_code=404, detail="File not found")

    try:
        # Use errors="replace" to ensure we always return a string even if bytes are not valid for the specified encoding.
        with file_path.open("r", encoding=encoding, errors="replace") as f:
            content = f.read()
    except LookupError:
        # Unknown encoding provided by client
        raise HTTPException(status_code=400, detail="Unknown encoding")
    except PermissionError:
        raise HTTPException(status_code=403, detail="Permission denied")

    rel_path = str(file_path.relative_to(user_base)).replace(os.sep, "/")
    return ReadFileResponse(path=rel_path, content=content, encoding=encoding)


class XlsxResponse(BaseModel):
    path: str
    sheetNames: list[str]
    sheet: str
    rows: list[list[str]]


class DocxBlock(BaseModel):
    type: Literal["paragraph", "table"]
    text: str | None = None
    rows: list[list[str]] | None = None


class DocxResponse(BaseModel):
    path: str
    blocks: list[DocxBlock]


@app.get("/api/docx", response_model=DocxResponse, response_model_exclude_none=True)
def read_docx(
    request: Request,
    path: str = Query(..., description="Relative .docx file path under user's base directory"),
):
    if DocxDocument is None or DocxParagraph is None or DocxTable is None or CT_P is None or CT_Tbl is None:
        raise HTTPException(status_code=500, detail="python-docx is not installed on the server")
    _, user_base = require_user(request)
    file_path = safe_join(user_base, path)
    if not file_path.exists() or not file_path.is_file():
        raise HTTPException(status_code=404, detail="File not found")
    if file_path.suffix.lower() != ".docx":
        raise HTTPException(status_code=400, detail="Unsupported Word format (only .docx)")

    try:
        document = DocxDocument(str(file_path))
    except PermissionError:
        raise HTTPException(status_code=403, detail="Permission denied")
    except Exception:
        raise HTTPException(status_code=400, detail="Failed to open document")

    blocks: list[DocxBlock] = []
    for item in _iter_docx_blocks(document):
        if isinstance(item, DocxParagraph):
            text = _docx_paragraph_to_text(item)
            if text:
                blocks.append(DocxBlock(type="paragraph", text=text))
        elif isinstance(item, DocxTable):
            table_rows: list[list[str]] = []
            for row in item.rows:
                row_cells = [_docx_cell_to_text(cell) for cell in row.cells]
                table_rows.append(row_cells)
            filtered_rows = [r for r in table_rows if any(cell for cell in r)]
            if filtered_rows:
                blocks.append(DocxBlock(type="table", rows=filtered_rows))

    rel_path = str(file_path.relative_to(user_base)).replace(os.sep, "/")
    return DocxResponse(path=rel_path, blocks=blocks)


@app.get("/api/xlsx", response_model=XlsxResponse)
def read_xlsx(
    request: Request,
    path: str = Query(..., description="Relative .xlsx/.xlsm file path under user's base directory"),
    sheet: str | None = Query(None, description="Sheet name to preview (defaults to the first sheet)"),
    maxRows: int = Query(500, ge=1, le=5000, description="Maximum number of rows to return"),
    maxCols: int = Query(100, ge=1, le=512, description="Maximum number of columns to return"),
):
    """Read an Excel workbook and return a table-like preview.

    Notes:
    - Requires openpyxl.
    - Returns values-only; formats are not preserved.
    - Datetimes and other types are converted to string.
    """
    if openpyxl is None:
        raise HTTPException(status_code=500, detail="openpyxl is not installed on the server")

    _, user_base = require_user(request)
    file_path = safe_join(user_base, path)
    if not file_path.exists() or not file_path.is_file():
        raise HTTPException(status_code=404, detail="File not found")
    if file_path.suffix.lower() not in {".xlsx", ".xlsm"}:
        raise HTTPException(status_code=400, detail="Unsupported Excel format (only .xlsx/.xlsm)")

    try:
        wb = openpyxl.load_workbook(filename=str(file_path), read_only=True, data_only=True)
    except PermissionError:
        raise HTTPException(status_code=403, detail="Permission denied")
    except Exception:
        raise HTTPException(status_code=400, detail="Failed to open workbook")

    sheet_names = wb.sheetnames
    if not sheet_names:
        return XlsxResponse(path=str(file_path.relative_to(BASE_DIR)).replace(os.sep, "/"), sheetNames=[], sheet="", rows=[])

    target_name = sheet if sheet in sheet_names else sheet_names[0]
    ws = wb[target_name]

    rows: list[list[str]] = []
    fetched_rows = 0
    for r in ws.iter_rows(min_row=1, max_row=maxRows, max_col=maxCols, values_only=True):
        row_out: list[str] = []
        for v in r:
            if v is None:
                row_out.append("")
            else:
                # Convert common types to string; avoid scientific notation surprises
                if isinstance(v, float):
                    # Trim trailing zeros while keeping a dot when needed
                    s = ("%f" % v).rstrip("0").rstrip(".")
                    row_out.append(s)
                else:
                    row_out.append(str(v))
        rows.append(row_out)
        fetched_rows += 1
        if fetched_rows >= maxRows:
            break

    rel = str(file_path.relative_to(user_base)).replace(os.sep, "/")
    return XlsxResponse(path=rel, sheetNames=sheet_names, sheet=target_name, rows=rows)


class XlsxWriteRequest(BaseModel):
    path: str
    sheet: str
    rows: list[list[str]]


@app.post("/api/xlsx/write", response_model=Entry)
def write_xlsx(request: Request, req: XlsxWriteRequest):
    if openpyxl is None:
        raise HTTPException(status_code=500, detail="openpyxl is not installed on the server")

    _, user_base = require_user(request)
    file_path = safe_join(user_base, req.path)
    if not file_path.exists() or not file_path.is_file():
        raise HTTPException(status_code=404, detail="File not found")
    if file_path.suffix.lower() not in {".xlsx", ".xlsm"}:
        raise HTTPException(status_code=400, detail="Unsupported Excel format (only .xlsx/.xlsm)")

    try:
        wb = openpyxl.load_workbook(filename=str(file_path))
    except PermissionError:
        raise HTTPException(status_code=403, detail="Permission denied")
    except Exception:
        raise HTTPException(status_code=400, detail="Failed to open workbook")

    if req.sheet not in wb.sheetnames:
        raise HTTPException(status_code=404, detail="Sheet not found")

    ws = wb[req.sheet]
    # Clear existing content
    if ws.max_row > 0:
        ws.delete_rows(1, ws.max_row)
    # Append new rows
    for row in req.rows:
        # Ensure values are simple (strings)
        ws.append([None if (v is None or str(v) == "") else str(v) for v in row])

    try:
        wb.save(str(file_path))
    except PermissionError:
        raise HTTPException(status_code=403, detail="Permission denied")
    except Exception:
        raise HTTPException(status_code=500, detail="Failed to save workbook")

    return stat_entry(file_path, user_base)


class WriteFileRequest(BaseModel):
    path: str
    content: str
    encoding: str = "utf-8"


@app.post("/api/write", response_model=Entry)
def write_file(request: Request, req: WriteFileRequest):
    _, user_base = require_user(request)
    file_path = safe_join(user_base, req.path)
    if not file_path.exists() or not file_path.is_file():
        raise HTTPException(status_code=404, detail="File not found")
    try:
        with file_path.open("w", encoding=req.encoding, newline="") as f:
            f.write(req.content)
    except LookupError:
        # Unknown encoding provided by client
        raise HTTPException(status_code=400, detail="Unknown encoding")
    except PermissionError:
        raise HTTPException(status_code=403, detail="Permission denied")
    return stat_entry(file_path, user_base)


class CreateFileRequest(BaseModel):
    parent: str = ""
    name: str
    content: str = ""
    encoding: str = "utf-8"


@app.post("/api/create", response_model=Entry)
def create_file(request: Request, req: CreateFileRequest):
    validate_name(req.name)
    _, user_base = require_user(request)
    target_dir = safe_join(user_base, req.parent)
    if not target_dir.exists() or not target_dir.is_dir():
        raise HTTPException(status_code=404, detail="Parent directory not found")
    dest = target_dir / req.name
    if dest.exists():
        raise HTTPException(status_code=409, detail="Already exists")
    try:
        with dest.open("w", encoding=req.encoding, newline="") as f:
            f.write(req.content)
    except LookupError:
        raise HTTPException(status_code=400, detail="Unknown encoding")
    except PermissionError:
        raise HTTPException(status_code=403, detail="Permission denied")
    return stat_entry(dest, user_base)


class MkdirRequest(BaseModel):
    parent: str = ""
    name: str


@app.post("/api/mkdir", response_model=Entry)
def mkdir(request: Request, req: MkdirRequest):
    validate_name(req.name)
    _, user_base = require_user(request)
    target_dir = safe_join(user_base, req.parent)
    if not target_dir.exists() or not target_dir.is_dir():
        raise HTTPException(status_code=404, detail="Parent directory not found")
    new_path = target_dir / req.name
    if new_path.exists():
        raise HTTPException(status_code=409, detail="Already exists")
    try:
        new_path.mkdir()
    except PermissionError:
        raise HTTPException(status_code=403, detail="Permission denied")
    return stat_entry(new_path, user_base)


@app.post("/api/upload", response_model=Entry)
async def upload(
    request: Request,
    background: BackgroundTasks,
    path: str = Form(""),
    file: UploadFile = File(...),
    source_url: str | None = Form(None),
):
    username, user_base = require_user(request)
    parent = safe_join(user_base, path)
    if not parent.exists() or not parent.is_dir():
        raise HTTPException(status_code=404, detail="Parent directory not found")
    filename = file.filename or ""
    validate_name(filename)
    dest = parent / filename
    if dest.exists():
        raise HTTPException(status_code=409, detail="File already exists")
    try:
        with dest.open("wb") as f:
            content = await file.read()
            f.write(content)
    except PermissionError:
        raise HTTPException(status_code=403, detail="Permission denied")

    suffix = dest.suffix.lower()
    if suffix in {".pdf", ".pptx", ".ppt"}:
        target_folder = dest.parent / dest.stem
        try:
            target_folder.mkdir(parents=True, exist_ok=True)
        except Exception:
            logger.exception("Failed to create folder %s for uploaded file %s", target_folder, dest.name)
            raise HTTPException(status_code=500, detail="ドキュメント用フォルダの作成に失敗しました")
        new_dest = target_folder / dest.name
        try:
            if new_dest.exists():
                if new_dest.is_file() or new_dest.is_symlink():
                    new_dest.unlink()
                else:
                    shutil.rmtree(new_dest)
            dest.rename(new_dest)
            dest = new_dest
        except PermissionError:
            raise HTTPException(status_code=403, detail="ファイルの配置に失敗しました")
        except Exception:
            logger.exception("Failed to relocate uploaded document %s into %s", dest.name, target_folder)
            raise HTTPException(status_code=500, detail="ファイルの配置に失敗しました")

        if suffix == ".pdf":
            # Initialize progress
            rel_after = dest.relative_to(user_base).as_posix()
            _set_progress(username, rel_after, stage="queued", percent=0.0, message="解析待機中", done=False)
            background.add_task(_process_pdf_upload, username, str(dest), source_url, None, dest.relative_to(user_base).as_posix())
        elif suffix in {".pptx", ".ppt"}:
            rel_after = dest.relative_to(user_base).as_posix()
            _set_progress(username, rel_after, stage="queued", percent=0.0, message="解析待機中", done=False)
            background.add_task(_process_pptx_upload, username, str(dest))

    return stat_entry(dest, user_base)


@app.get("/api/upload/status")
def get_upload_status(request: Request, path: str = Query(...)) -> dict:
    """Return progress for a processed document path (relative to user base)."""
    username, user_base = require_user(request)
    file_path = safe_join(user_base, path)
    try:
        rel = file_path.relative_to(user_base).as_posix()
    except Exception:
        raise HTTPException(status_code=400, detail="Invalid path")
    data = _get_progress(username, rel)
    if not data:
        if file_path.exists():
            return {"stage": "done", "percent": 1.0, "message": "完了", "done": True}
        return {"stage": "unknown", "percent": 0.0, "message": "未開始", "done": False}
    return data


@app.delete("/api/delete")
def delete(request: Request, path: str = Query(..., description="Relative path to file or directory under user's base directory")):
    # Important: do not resolve the final leaf to ensure we unlink symlinks rather than deleting targets
    _, user_base = require_user(request)
    target = safe_join_leaf(user_base, path)
    if not target.exists():
        raise HTTPException(status_code=404, detail="Path not found")
    try:
        # Never follow symlinks when deleting: unlink the link itself
        if target.is_symlink():
            target.unlink()
        elif target.is_dir():
            shutil.rmtree(target)
        else:
            target.unlink()
    except PermissionError:
        raise HTTPException(status_code=403, detail="Permission denied")
    return {"ok": True}


class RenameRequest(BaseModel):
    path: str
    newName: str


@app.post("/api/rename", response_model=Entry)
def rename(request: Request, req: RenameRequest):
    validate_name(req.newName)
    _, user_base = require_user(request)
    src = safe_join(user_base, req.path)
    if not src.exists():
        raise HTTPException(status_code=404, detail="Path not found")
    dst = src.with_name(req.newName)
    if dst.exists():
        raise HTTPException(status_code=409, detail="Target already exists")
    try:
        src.rename(dst)
    except PermissionError:
        raise HTTPException(status_code=403, detail="Permission denied")
    return stat_entry(dst, user_base)


class MoveRequest(BaseModel):
    src: str
    dstDir: str
    newName: str | None = None


@app.post("/api/move", response_model=Entry)
def move(request: Request, req: MoveRequest):
    """Move a file or directory to a destination directory (optionally with a new name)."""
    _, user_base = require_user(request)
    src = safe_join(user_base, req.src)
    if not src.exists():
        raise HTTPException(status_code=404, detail="Source not found")
    dst_dir = safe_join(user_base, req.dstDir)
    if not dst_dir.exists() or not dst_dir.is_dir():
        raise HTTPException(status_code=404, detail="Destination directory not found")

    target_name = req.newName or src.name
    validate_name(target_name)
    dst = (dst_dir / target_name)
    if dst.exists():
        raise HTTPException(status_code=409, detail="Target already exists")

    try:
        # Prevent moving a directory into its own subdirectory
        if src.is_dir():
            try:
                (dst_dir / "").resolve().relative_to(src.resolve())
                raise HTTPException(status_code=400, detail="Cannot move a directory into its own subdirectory")
            except ValueError:
                pass
        # shutil.move handles cross-device moves and both files/directories
        shutil.move(str(src), str(dst))
    except PermissionError:
        raise HTTPException(status_code=403, detail="Permission denied")
    return stat_entry(Path(dst), user_base)


class CopyRequest(BaseModel):
    src: str
    dstDir: str
    newName: str | None = None


@app.post("/api/copy", response_model=Entry)
def copy(request: Request, req: CopyRequest):
    """Copy a file or directory to a destination directory (optionally with a new name)."""
    _, user_base = require_user(request)
    src = safe_join(user_base, req.src)
    if not src.exists():
        raise HTTPException(status_code=404, detail="Source not found")
    dst_dir = safe_join(user_base, req.dstDir)
    if not dst_dir.exists() or not dst_dir.is_dir():
        raise HTTPException(status_code=404, detail="Destination directory not found")

    target_name = req.newName or src.name
    validate_name(target_name)
    dst = (dst_dir / target_name)
    if dst.exists():
        raise HTTPException(status_code=409, detail="Target already exists")

    try:
        if src.is_dir():
            # Prevent copying a directory into its own subdirectory
            try:
                (dst_dir / "").resolve().relative_to(src.resolve())
                raise HTTPException(status_code=400, detail="Cannot copy a directory into its own subdirectory")
            except ValueError:
                pass
            shutil.copytree(src, dst)
        else:
            shutil.copy2(src, dst)
    except PermissionError:
        raise HTTPException(status_code=403, detail="Permission denied")
    return stat_entry(Path(dst), user_base)


class SymlinkRequest(BaseModel):
    target: str
    dstDir: str
    name: str | None = None


@app.post("/api/symlink", response_model=Entry)
def create_symlink(request: Request, req: SymlinkRequest):
    """Create a symbolic link under dstDir pointing to target (both under BASE_DIR).

    - Rejects if target does not exist
    - Rejects if dstDir is not a directory
    - Rejects if the resulting link path already exists
    - Prevents linking a directory into its own subdirectory (would create cycles)
    """
    _, user_base = require_user(request)
    src = safe_join(user_base, req.target)
    if not src.exists():
        raise HTTPException(status_code=404, detail="Source not found")
    dst_dir = safe_join(user_base, req.dstDir)
    if not dst_dir.exists() or not dst_dir.is_dir():
        raise HTTPException(status_code=404, detail="Destination directory not found")

    link_name = (req.name or src.name)
    validate_name(link_name)
    link_path = (dst_dir / link_name)
    if link_path.exists():
        raise HTTPException(status_code=409, detail="Target already exists")

    try:
        # Prevent linking a directory into its own subdirectory
        if src.is_dir():
            try:
                (dst_dir / "").resolve().relative_to(src.resolve())
                raise HTTPException(status_code=400, detail="Cannot link a directory into its own subdirectory")
            except ValueError:
                pass

        # Create a relative symlink to keep tree relocatable
        rel_target = os.path.relpath(str(src), str(dst_dir))
        # On POSIX, target_is_directory is ignored; on Windows it's required
        link_path.symlink_to(rel_target, target_is_directory=src.is_dir())
    except PermissionError:
        raise HTTPException(status_code=403, detail="Permission denied")
    except NotImplementedError:
        raise HTTPException(status_code=500, detail="Symlinks not supported on this platform")

    return stat_entry(link_path, user_base)


@app.get("/api/info")
def info(request: Request):
    username, user_base = require_user(request)
    return {
        "baseDir": str(user_base),
        "workspace": str(WORKSPACE_ROOT),
        "serverTime": datetime.utcnow().isoformat() + "Z",
        "username": username,
    }


# --------------- Agent Chat (OpenAI/Azure + MCP tools) -----------------

class AgentMessage(BaseModel):
    role: Literal["system", "user", "assistant"]
    content: str


class AgentChatRequest(BaseModel):
    folder: str = ""  # relative path under BASE_DIR; empty means database root
    messages: list[AgentMessage]
    maxFiles: int = 5


class AgentChatResponse(BaseModel):
    messages: list[AgentMessage]
    usedFiles: list[str]


class AgentFileChatRequest(BaseModel):
    file: str  # relative path under BASE_DIR
    messages: list[AgentMessage]


def _http_get_json(url: str, authorization: str | None = None) -> dict | list:
    headers = {"accept": "application/json"}
    if authorization:
        headers["authorization"] = authorization
    req = urllib.request.Request(url, headers=headers)
    with urllib.request.urlopen(req, timeout=30) as resp:  # nosec B310
        data = resp.read()
        ct = resp.headers.get("content-type", "")
        if "application/json" not in ct:
            # Best-effort parse; raise if fails
            try:
                return json.loads(data.decode("utf-8", errors="replace"))
            except Exception:
                raise HTTPException(status_code=502, detail=f"Unexpected content-type: {ct or 'unknown'}")
        return json.loads(data.decode("utf-8", errors="replace"))


def _http_post_json(url: str, body: dict, authorization: str | None = None) -> dict | list:
    headers = {"accept": "application/json", "content-type": "application/json"}
    if authorization:
        headers["authorization"] = authorization
    data = json.dumps(body, ensure_ascii=False).encode("utf-8")
    req = urllib.request.Request(url, headers=headers, data=data, method="POST")
    with urllib.request.urlopen(req, timeout=30) as resp:  # nosec B310
        raw = resp.read()
        ct = resp.headers.get("content-type", "")
        if "application/json" not in ct:
            try:
                return json.loads(raw.decode("utf-8", errors="replace"))
            except Exception:
                raise HTTPException(status_code=502, detail=f"Unexpected content-type: {ct or 'unknown'}")
        return json.loads(raw.decode("utf-8", errors="replace"))


def _http_delete_json(url: str, authorization: str | None = None) -> dict | list:
    headers = {"accept": "application/json"}
    if authorization:
        headers["authorization"] = authorization
    req = urllib.request.Request(url, headers=headers, method="DELETE")
    with urllib.request.urlopen(req, timeout=30) as resp:  # nosec B310
        raw = resp.read()
        ct = resp.headers.get("content-type", "")
        if not raw:
            return {"ok": True}
        if "application/json" not in ct:
            try:
                return json.loads(raw.decode("utf-8", errors="replace"))
            except Exception:
                raise HTTPException(status_code=502, detail=f"Unexpected content-type: {ct or 'unknown'}")
        return json.loads(raw.decode("utf-8", errors="replace"))


def _tool_protocol_instructions() -> str:
    return (
        "You are an assistant with access to MCP tools.\n"
        "Available tools: get_meta(path), read_file(path), create_file(parent,name,content,encoding='utf-8'), "
        "update_file(path,content,encoding='utf-8'), delete_file(path).\n"
        "When you want to use a tool, respond ONLY with a single JSON object on one line: "
        "{\"action\":\"tool\",\"tool\":\"<name>\",\"args\":{...}}.\n"
        "When you are ready to answer the user, respond ONLY with {\"action\":\"final\",\"content\":\"...\"}.\n"
        "Paths must be within the user's directory. For safety, prefer absolute MCP paths like '<username>/path/under/root'.\n"
        "For any creation or update request (e.g., \"markdownを作成\"), you MUST actually call the corresponding tool (create_file/update_file). Do not claim success in the final answer without using tools.\n"
    )


# Normalize a tool path to include the username and optionally the selected base folder
def _normalize_mcp_path(raw_path: str, username: str, base_folder_rel: str | None) -> str:
    p = (raw_path or "").replace("\\", "/").strip("/")
    u = (username or "").strip().strip("/")
    base = (base_folder_rel or "").replace("\\", "/").strip("/")
    if not p:
        return u
    # Already absolute under user
    if p.lower().startswith(f"{u.lower()}/") or p.lower() == u.lower():
        return p
    # If base folder known, treat p as relative to it
    if base:
        # Avoid double-prefix if p already starts with base
        if p.lower().startswith(base.lower() + "/") or p.lower() == base.lower():
            return "/".join([u, p])
        return "/".join([u, base, p])
    # Otherwise treat as relative to user root
    return "/".join([u, p])

# Extract a tool directive JSON object from arbitrary model text output.
# Accepts raw text, code-fenced JSON (```json ... ```), or multi-line content containing a JSON object.
def _extract_json_directive(text: str) -> dict | None:
    try:
        s = (text or "").strip()
        # Strip code fences if present
        if s.startswith("```"):
            # Remove first fence line and trailing fence
            lines = s.splitlines()
            # drop first line like ```json or ```
            if lines:
                lines = lines[1:]
            # remove trailing ``` if present
            if lines and lines[-1].strip().startswith("```"):
                lines = lines[:-1]
            s = "\n".join(lines).strip()
        # Fast path: whole content is a JSON object
        try:
            obj = json.loads(s)
            if isinstance(obj, dict) and obj.get("action") in ("tool", "final"):
                return obj
        except Exception:
            pass
        # Fallback: find first {...} block heuristically
        start = s.find("{")
        while start != -1:
            # attempt to find a matching closing brace by incremental expansion
            end = start + 1
            depth = 1
            while end < len(s) and depth > 0:
                ch = s[end]
                if ch == '{':
                    depth += 1
                elif ch == '}':
                    depth -= 1
                end += 1
            if depth == 0:
                segment = s[start:end]
                try:
                    obj = json.loads(segment)
                    if isinstance(obj, dict) and obj.get("action") in ("tool", "final"):
                        return obj
                except Exception:
                    pass
            start = s.find("{", start + 1)
    except Exception:
        return None
    return None

# Centralized non-streaming tool-calling loop used by multiple endpoints
def _run_tool_calling(messages: list[dict], profile: LLMProfile, auth_header: str) -> tuple[str, list[str]]:
    used_files: list[str] = []

    # If Agents SDK is available and OpenAI is preferred, run via Agents SDK with MCP server
    if Agent is not None and MCPServerStreamableHttp is not None and profile.preferred == "openai":
        try:
            cfg = profile.openai
            if not cfg.apiKey:
                raise RuntimeError("OpenAI API key is not configured")
            # Build MCP server connection (streamable HTTP) with auth header forwarding
            async def run_agent_sync() -> str:
                import asyncio
                async def _run() -> str:
                    # Configure environment for Agents SDK (API key/base URL)
                    prev_key = os.environ.get("OPENAI_API_KEY")
                    prev_base = os.environ.get("OPENAI_BASE_URL")
                    try:
                        os.environ["OPENAI_API_KEY"] = cfg.apiKey
                        if cfg.baseUrl:
                            os.environ["OPENAI_BASE_URL"] = cfg.baseUrl
                        model_name = (cfg.model or os.environ.get("OPENAI_MODEL") or "gpt-4o-mini")
                        async with MCPServerStreamableHttp(
                            name="Directory RAG MCP",
                            params={
                            "url": f"{MCP_SERVER_BASE}/mcp",
                            "headers": {"authorization": auth_header},
                            "timeout": 15,
                        },
                        cache_tools_list=True,
                        max_retry_attempts=2,
                    ) as server:
                            system_texts = [m["content"] for m in messages if m.get("role") == "system"]
                            instructions = "\n\n".join(system_texts) if system_texts else _load_system_prompt()
                            user_messages = [m for m in messages if m.get("role") != "system"]
                            agent = Agent(
                                name="Assistant",
                                instructions=instructions,
                                mcp_servers=[server],
                                model=model_name,
                                model_settings=ModelSettings(tool_choice="auto"),
                            )
                            # Convert to agents SDK format: pass concatenated messages as a single input
                            # Keep it simple: combine user/assistant turns into a plain transcript
                            transcript = []
                            for m in user_messages:
                                role = m.get("role") or "user"
                                content = m.get("content") or ""
                                prefix = "User:" if role == "user" else "Assistant:"
                                transcript.append(f"{prefix} {content}")
                            prompt = "\n".join(transcript)
                            result = await Runner.run(agent, prompt)
                            return result.final_output or ""
                    finally:
                        # Restore environment
                        if prev_key is not None:
                            os.environ["OPENAI_API_KEY"] = prev_key
                        else:
                            os.environ.pop("OPENAI_API_KEY", None)
                        if prev_base is not None:
                            os.environ["OPENAI_BASE_URL"] = prev_base
                        else:
                            os.environ.pop("OPENAI_BASE_URL", None)
                return asyncio.run(_run())

            final = run_agent_sync()
            return final, used_files
        except Exception as e:
            # Fall back to the legacy loop below
            messages.append({"role": "system", "content": f"[agents-sdk-fallback] {e}"})

    # Legacy manual tool-calling loop (Azure or Agents SDK unavailable)
    def _complete_directive(msgs: list[dict]) -> tuple[str, dict]:
        try:
            if profile.preferred == "azure":
                cfg = profile.azure
                if not (cfg.apiKey and cfg.endpoint and cfg.deployment):
                    raise RuntimeError("Azure OpenAI is not fully configured")
                client = AzureOpenAI(azure_endpoint=cfg.endpoint, api_key=cfg.apiKey, api_version=(cfg.apiVersion or "2024-02-15-preview"))
                comp = client.chat.completions.create(model=cfg.deployment, messages=msgs, temperature=0.1)
            else:
                cfg = profile.openai
                if not cfg.apiKey:
                    raise RuntimeError("OpenAI API key is not configured")
                client = OpenAI(api_key=cfg.apiKey, base_url=(cfg.baseUrl or None))
                model_name = (cfg.model or os.environ.get("OPENAI_MODEL") or "gpt-4o-mini")
                comp = client.chat.completions.create(model=model_name, messages=msgs, temperature=0.1)
            text = (comp.choices[0].message.content or "").strip()
            obj = _extract_json_directive(text)
            if isinstance(obj, dict) and obj.get("action") in ("tool", "final"):
                return "ok", obj
            return "parse_error", {"raw": text}
        except Exception as ex:
            return "error", {"message": f"{ex}"}

    max_steps = 8
    step = 0
    last_answer = ""
    while step < max_steps:
        step += 1
        status, obj = _complete_directive(messages + [{"role": "system", "content": "Respond with a single JSON object as instructed."}])
        if status == "ok" and isinstance(obj, dict):
            if obj.get("action") == "final":
                last_answer = str(obj.get("content") or "")
                break
            if obj.get("action") == "tool":
                tool = str(obj.get("tool") or "").strip()
                args = obj.get("args") or {}
                try:
                    if tool == "get_meta":
                        p = str(args.get("path") or "")
                        url = f"{MCP_SERVER_BASE}/tools/get_meta?" + urllib.parse.urlencode({"path": p})
                        result = _http_get_json(url, auth_header)
                    elif tool == "read_file":
                        p = str(args.get("path") or "")
                        url = f"{MCP_SERVER_BASE}/tools/read_file?" + urllib.parse.urlencode({"path": p})
                        result = _http_get_json(url, auth_header)
                        if isinstance(p, str) and p:
                            try:
                                used_files.append(p)
                            except Exception:
                                pass
                    elif tool == "create_file":
                        parent = str(args.get("parent") or "")
                        name = str(args.get("name") or "")
                        content = str(args.get("content") or "")
                        encoding = str(args.get("encoding") or "utf-8")
                        url = f"{MCP_SERVER_BASE}/tools/create_file"
                        result = _http_post_json(url, {"parent": parent, "name": name, "content": content, "encoding": encoding}, auth_header)
                    elif tool == "update_file":
                        p = str(args.get("path") or "")
                        content = str(args.get("content") or "")
                        encoding = str(args.get("encoding") or "utf-8")
                        url = f"{MCP_SERVER_BASE}/tools/update_file"
                        result = _http_post_json(url, {"path": p, "content": content, "encoding": encoding}, auth_header)
                    elif tool == "delete_file":
                        p = str(args.get("path") or "")
                        url = f"{MCP_SERVER_BASE}/tools/delete_file?" + urllib.parse.urlencode({"path": p})
                        result = _http_delete_json(url, auth_header)
                    else:
                        result = {"error": f"Unknown tool: {tool}"}
                    try:
                        snippet = json.dumps(result, ensure_ascii=False) if not isinstance(result, str) else result
                    except Exception:
                        snippet = str(result)
                    messages.append({"role": "system", "content": f"Tool {tool} result (JSON):\n{snippet}"})
                    continue
                except Exception as tex:
                    messages.append({"role": "system", "content": f"Tool {tool} error: {tex}"})
                    continue
        break

    if not last_answer:
        try:
            if profile.preferred == "azure":
                cfg = profile.azure
                if not (cfg.apiKey and cfg.endpoint and cfg.deployment):
                    raise RuntimeError("Azure OpenAI is not fully configured")
                client = AzureOpenAI(azure_endpoint=cfg.endpoint, api_key=cfg.apiKey, api_version=(cfg.apiVersion or "2024-02-15-preview"))
                comp = client.chat.completions.create(model=cfg.deployment, messages=messages, temperature=0.2)  # type: ignore[arg-type]
            else:
                cfg = profile.openai
                if not cfg.apiKey:
                    raise RuntimeError("OpenAI API key is not configured")
                client = OpenAI(api_key=cfg.apiKey, base_url=(cfg.baseUrl or None))
                model_name = (cfg.model or os.environ.get("OPENAI_MODEL") or "gpt-4o-mini")
                comp = client.chat.completions.create(model=model_name, messages=messages, temperature=0.2)  # type: ignore[arg-type]
            last_answer = (comp.choices[0].message.content or "")
        except Exception as e:
            last_answer = f"[Error: {e}]"

    return last_answer, used_files


def _load_system_prompt() -> str:
    prompt_path = BASE_DIR / "agent_system_prompt.md"
    if prompt_path.exists() and prompt_path.is_file():
        try:
            return prompt_path.read_text(encoding="utf-8", errors="replace")
        except Exception:
            pass
    # Fallback minimal prompt
    return (
        "You are a file-aware assistant. Use MCP tools to browse and modify files when helpful. "
        + _tool_protocol_instructions()
    )


def _flatten_tree(base_rel: str, node: dict, prefix: str = "", is_root: bool = True) -> list[str]:
    """Flatten DirNode/FileNode into relative file paths under the database root.

    The returned paths are relative to the database root. If a subfolder `base_rel`
    is selected, file paths are prefixed by that folder but the root node name from
    MCP (`name` of the selected folder) is not duplicated.
    """
    paths: list[str] = []
    node_type = node.get("type")
    node_name = node.get("name", "")

    if node_type == "file":
        parts = []  # build: [base_rel?, prefix?, file]
        if base_rel:
            parts.append(base_rel)
        if prefix:
            parts.append(prefix)
        parts.append(node_name)
        rel = "/".join([p for p in parts if p]).strip("/")
        paths.append(rel)
        return paths

    if node_type == "folder":
        # For the root folder node, do not include its name in the prefix.
        # For nested folders, include their name.
        next_prefix = prefix
        if not is_root:
            next_prefix = f"{prefix}/{node_name}".strip("/") if prefix else node_name
        for child in (node.get("children") or []):
            paths.extend(_flatten_tree(base_rel, child, next_prefix, False))
        return paths

    return paths


def _is_identity_or_smalltalk(question: str) -> bool:
    q = (question or "").lower()
    # If the question clearly mentions files/folders, do not treat as smalltalk
    file_words_jp = ["ファイル", "フォルダ", "ディレクトリ", "一覧", "見せて", "開いて", "内容", "中身", "パス"]
    file_words_en = ["file", "files", "folder", "directory", "list", "show", "open", "content", "contents", "path"]
    if any(w in question for w in file_words_jp) or any(w in q for w in file_words_en):
        return False

    jp_triggers = [
        "あなたは誰",
        "誰ですか",
        "自己紹介",
        "何者",
        "何ができる",
        "できること",
        "どんな存在",
        "どんなai",
        "よろしく",
        "はじめまして",
        "こんにちは",
        "こんばんは",
        "おはよう",
        "名前は",
        "君は誰",
    ]
    en_triggers = [
        "who are you",
        "what are you",
        "introduce yourself",
        "about you",
        "your name",
        "what can you do",
        "capabilities",
        "hello",
        "hi",
        "hey",
        "greetings",
    ]
    if any(t in question for t in jp_triggers):
        return True
    if any(t in q for t in en_triggers):
        return True
    return False


def _select_candidate_files(question: str, all_files: list[str], max_files: int) -> list[str]:
    # For identity/smalltalk queries, skip file reads entirely
    if _is_identity_or_smalltalk(question):
        return []

    question_lower = question.lower()
    keywords = {w for w in re.split(r"[^a-zA-Z0-9_.-]+", question_lower) if w}
    scored: list[tuple[int, str]] = []
    for p in all_files:
        name = p.lower()
        score = 0
        for k in keywords:
            if k and k in name:
                score += 2
        # prefer markdown/csv/txt/json
        if name.endswith((".md", ".markdown", ".mdx")):
            score += 3
        if name.endswith((".txt", ".log", ".json")):
            score += 2
        if name.endswith((".csv", ".tsv", ".xlsx", ".xlsm")):
            score += 1
        scored.append((score, p))
    scored.sort(key=lambda t: (-t[0], t[1]))
    # Only take positively scored items; allow zero candidates
    top = [p for s, p in scored if s > 0][: max(0, max_files)]
    return top


def _read_single_file_for_chat(username: str, file_rel: str, request: Request) -> tuple[str, str]:
    normalized = file_rel.replace("\\", "/").strip("/")
    if not normalized:
        raise HTTPException(status_code=400, detail="Invalid file path")
    mcp_path = "/".join(part for part in [username, normalized] if part)
    rf_url = f"{MCP_SERVER_BASE}/tools/read_file?{urllib.parse.urlencode({'path': mcp_path})}"
    data = _http_get_json(rf_url, _get_auth_header(request))
    content: str
    if isinstance(data, dict) and "content" in data and isinstance(data["content"], str):
        content = data["content"]
    elif isinstance(data, dict) and all(k in data for k in ("sheetNames", "sheet", "rows")):
        content = json.dumps(data, ensure_ascii=False)
    else:
        content = json.dumps(data, ensure_ascii=False)
    if len(content) > 20000:
        content = content[:20000] + "\n...[truncated]"
    return mcp_path, content


@app.post("/api/agent/chat", response_model=AgentChatResponse)
def agent_chat(request: Request, req: AgentChatRequest):
    if not req.messages or req.messages[-1].role != "user":
        raise HTTPException(status_code=400, detail="Last message must be from user")

    username, user_base = require_user(request)
    folder_rel = req.folder.replace("\\", "/").strip("/")
    target = safe_join(user_base, folder_rel)
    if not target.exists() or not target.is_dir():
        raise HTTPException(status_code=404, detail="Target folder not found")

    system_prompt = _load_system_prompt()
    messages: list[dict] = [{"role": "system", "content": system_prompt}]
    messages.extend({"role": m.role, "content": m.content} for m in req.messages)
    # Provide explicit MCP folder path hint so the model can call tools
    mcp_folder = f"{username}/{folder_rel}" if folder_rel else username
    messages.append({"role": "system", "content": f"Selected folder MCP path: {mcp_folder}. Use get_meta(path) to browse and read_file(path) if needed."})
    messages.append({"role": "system", "content": _tool_protocol_instructions()})

    profile = _load_llm_profile(username)
    answer, used_files = _run_tool_calling(messages, profile, _get_auth_header(request))
    out_messages = list(req.messages)
    out_messages.append(AgentMessage(role="assistant", content=answer or ""))
    return AgentChatResponse(messages=out_messages, usedFiles=used_files)


@app.post("/api/agent/chat/stream")
def agent_chat_stream(request: Request, req: AgentChatRequest):
    if not req.messages or req.messages[-1].role != "user":
        raise HTTPException(status_code=400, detail="Last message must be from user")

    username, user_base = require_user(request)
    folder_rel = req.folder.replace("\\", "/").strip("/")
    target = safe_join(user_base, folder_rel)
    if not target.exists() or not target.is_dir():
        raise HTTPException(status_code=404, detail="Target folder not found")

    # Defer MCP calls and message construction into the stream to enable real-time tool status
    def ndjson_iter():
        mcp_folder = f"{username}/{folder_rel}" if folder_rel else username
        used_files: list[str] = []
        # Build initial messages; delegate tool decisions to the LLM
        system_prompt = _load_system_prompt()
        messages: list[dict] = [{"role": "system", "content": system_prompt}]
        for m in req.messages:
            messages.append({"role": m.role, "content": m.content})
        # Provide a light hint on selected folder scope
        messages.append({"role": "system", "content": f"Selected folder MCP path: {mcp_folder}. When specifying tool paths, prefer absolute form '<username>/...'. If the model supplies a relative path, it will be resolved against the selected folder."})
        messages.append({"role": "system", "content": _tool_protocol_instructions()})
        # Emit initial meta (no used files yet)
        try:
            yield (json.dumps({"type": "meta", "usedFiles": used_files}, ensure_ascii=False) + "\n").encode("utf-8")
        except Exception:
            pass

        profile = _load_llm_profile(username)

        # Always use the interactive tool loop to surface tool events in streaming

        # Helper to get a single JSON directive from the model
        def _complete_directive(msgs: list[dict]) -> tuple[str, dict]:
            try:
                if profile.preferred == "azure":
                    cfg = profile.azure
                    if not (cfg.apiKey and cfg.endpoint and cfg.deployment):
                        raise RuntimeError("Azure OpenAI is not fully configured")
                    client = AzureOpenAI(azure_endpoint=cfg.endpoint, api_key=cfg.apiKey, api_version=(cfg.apiVersion or "2024-02-15-preview"))
                    comp = client.chat.completions.create(model=cfg.deployment, messages=msgs, temperature=0.1)
                else:
                    cfg = profile.openai
                    if not cfg.apiKey:
                        raise RuntimeError("OpenAI API key is not configured")
                    client = OpenAI(api_key=cfg.apiKey, base_url=(cfg.baseUrl or None))
                    model_name = (cfg.model or os.environ.get("OPENAI_MODEL") or "gpt-4o-mini")
                    comp = client.chat.completions.create(model=model_name, messages=msgs, temperature=0.1)
                text = (comp.choices[0].message.content or "").strip()
                obj = _extract_json_directive(text)
                if isinstance(obj, dict) and obj.get("action") in ("tool", "final"):
                    return "ok", obj
                return "parse_error", {"raw": text}
            except Exception as ex:
                return "error", {"message": f"{ex}"}

        # Interactive tool loop with bounded steps
        max_steps = 20
        step = 0
        last_answer = ""
        while step < max_steps:
            step += 1
            status, obj = _complete_directive(messages + [{"role": "system", "content": "Respond with a single JSON object as instructed."}])
            if status == "ok" and isinstance(obj, dict):
                if obj.get("action") == "final":
                    last_answer = str(obj.get("content") or "")
                    break
                if obj.get("action") == "tool":
                    tool = str(obj.get("tool") or "").strip()
                    args = obj.get("args") or {}
                    # Normalize path arguments to absolute MCP path under username
                    if "path" in args and isinstance(args.get("path"), str):
                        args["path"] = _normalize_mcp_path(str(args.get("path") or ""), username, folder_rel)
                    if "parent" in args and isinstance(args.get("parent"), str):
                        args["parent"] = _normalize_mcp_path(str(args.get("parent") or ""), username, folder_rel)
                    # Emit start event
                    try:
                        yield (json.dumps({"type": "tool", "tool": tool, "status": "start", "path": str(args.get("path") or args.get("parent") or "")}, ensure_ascii=False) + "\n").encode("utf-8")
                    except Exception:
                        pass
                    try:
                        auth = _get_auth_header(request)
                        result: dict | list | str
                        if tool == "get_meta":
                            p = str(args.get("path") or "")
                            url = f"{MCP_SERVER_BASE}/tools/get_meta?" + urllib.parse.urlencode({"path": p})
                            result = _http_get_json(url, auth)
                        elif tool == "read_file":
                            p = str(args.get("path") or "")
                            url = f"{MCP_SERVER_BASE}/tools/read_file?" + urllib.parse.urlencode({"path": p})
                            result = _http_get_json(url, auth)
                            # Track used file for UI
                            if isinstance(p, str) and p:
                                try:
                                    used_files.append(p)
                                except Exception:
                                    pass
                        elif tool == "create_file":
                            parent = str(args.get("parent") or "")
                            name = str(args.get("name") or "")
                            content = str(args.get("content") or "")
                            encoding = str(args.get("encoding") or "utf-8")
                            url = f"{MCP_SERVER_BASE}/tools/create_file"
                            result = _http_post_json(url, {"parent": parent, "name": name, "content": content, "encoding": encoding}, auth)
                        elif tool == "update_file":
                            p = str(args.get("path") or "")
                            content = str(args.get("content") or "")
                            encoding = str(args.get("encoding") or "utf-8")
                            url = f"{MCP_SERVER_BASE}/tools/update_file"
                            result = _http_post_json(url, {"path": p, "content": content, "encoding": encoding}, auth)
                        elif tool == "delete_file":
                            p = str(args.get("path") or "")
                            url = f"{MCP_SERVER_BASE}/tools/delete_file?" + urllib.parse.urlencode({"path": p})
                            result = _http_delete_json(url, auth)
                        else:
                            result = {"error": f"Unknown tool: {tool}"}
                        # Emit success (for create_file, include full created path if available)
                        try:
                            created_path = str(args.get("path") or args.get("parent") or "")
                            if tool == "create_file":
                                try:
                                    nm = str(args.get("name") or "").strip()
                                    if created_path and nm:
                                        created_path = "/".join(p for p in [created_path.strip("/"), nm] if p)
                                except Exception:
                                    pass
                            yield (json.dumps({"type": "tool", "tool": tool, "status": "success", "path": created_path}, ensure_ascii=False) + "\n").encode("utf-8")
                        except Exception:
                            pass
                        # Summarize tool result for the model
                        try:
                            snippet = json.dumps(result, ensure_ascii=False) if not isinstance(result, str) else result
                        except Exception:
                            snippet = str(result)
                        messages.append({"role": "system", "content": f"Tool {tool} result (JSON):\n{snippet}"})
                        continue
                    except Exception as tex:
                        # Emit error and append result for model to react
                        try:
                            yield (json.dumps({"type": "tool", "tool": tool, "status": "error", "message": f"{tex}"}, ensure_ascii=False) + "\n").encode("utf-8")
                        except Exception:
                            pass
                        messages.append({"role": "system", "content": f"Tool {tool} error: {tex}"})
                        continue
            # On any error or parse error, break and let model answer directly
            break

        # If we didn't get a final from the directive loop, do a normal completion
        if not last_answer:
            try:
                if profile.preferred == "azure":
                    cfg = profile.azure
                    if not (cfg.apiKey and cfg.endpoint and cfg.deployment):
                        raise RuntimeError("Azure OpenAI is not fully configured")
                    client = AzureOpenAI(azure_endpoint=cfg.endpoint, api_key=cfg.apiKey, api_version=(cfg.apiVersion or "2024-02-15-preview"))
                    stream = client.chat.completions.create(model=cfg.deployment, messages=messages, temperature=0.2, stream=True)
                else:
                    cfg = profile.openai
                    if not cfg.apiKey:
                        raise RuntimeError("OpenAI API key is not configured")
                    client = OpenAI(api_key=cfg.apiKey, base_url=(cfg.baseUrl or None))
                    model_name = (cfg.model or os.environ.get("OPENAI_MODEL") or "gpt-4o-mini")
                    stream = client.chat.completions.create(model=model_name, messages=messages, temperature=0.2, stream=True)
                # Stream chunks as deltas
                for chunk in stream:  # type: ignore[assignment]
                    try:
                        choice = (chunk.choices[0] if getattr(chunk, "choices", None) else None)
                        delta_obj = getattr(choice, "delta", None) if choice is not None else None
                        piece = getattr(delta_obj, "content", None) if delta_obj is not None else None
                        if isinstance(piece, str) and piece:
                            yield (json.dumps({"type": "delta", "content": piece}, ensure_ascii=False) + "\n").encode("utf-8")
                    except Exception:
                        # ignore malformed chunks
                        pass
                yield (json.dumps({"type": "done"}, ensure_ascii=False) + "\n").encode("utf-8")
                return
            except Exception as e:
                yield (json.dumps({"type": "error", "message": f"OpenAI error: {e}"}, ensure_ascii=False) + "\n").encode("utf-8")
                yield (json.dumps({"type": "done"}, ensure_ascii=False) + "\n").encode("utf-8")
                return

        if last_answer:
            # emit as a single delta
            yield (json.dumps({"type": "delta", "content": last_answer}, ensure_ascii=False) + "\n").encode("utf-8")
        yield (json.dumps({"type": "done"}, ensure_ascii=False) + "\n").encode("utf-8")

    return StreamingResponse(ndjson_iter(), media_type="application/x-ndjson")


@app.post("/api/agent/chat/file", response_model=AgentChatResponse)
def agent_chat_file(request: Request, req: AgentFileChatRequest):
    if not req.messages or req.messages[-1].role != "user":
        raise HTTPException(status_code=400, detail="Last message must be from user")

    username, user_base = require_user(request)
    file_rel = req.file.replace("\\", "/").strip("/")
    # Validate the path but do not pre-read; allow non-existent files (model may choose not to read)
    _ = safe_join(user_base, file_rel)
    mcp_path = "/".join(part for part in [username, file_rel] if part)

    system_prompt = _load_system_prompt()
    messages: list[dict] = [{"role": "system", "content": system_prompt}]
    messages.extend({"role": m.role, "content": m.content} for m in req.messages)
    messages.append({"role": "system", "content": _tool_protocol_instructions()})
    messages.append({"role": "system", "content": f"Selected file MCP path: {mcp_path}. Call read_file(path) only if needed."})

    profile = _load_llm_profile(username)
    answer, used_files = _run_tool_calling(messages, profile, _get_auth_header(request))
    out_messages = list(req.messages)
    out_messages.append(AgentMessage(role="assistant", content=answer or ""))
    return AgentChatResponse(messages=out_messages, usedFiles=used_files)


@app.post("/api/agent/chat/file/stream")
def agent_chat_file_stream(request: Request, req: AgentFileChatRequest):
    if not req.messages or req.messages[-1].role != "user":
        raise HTTPException(status_code=400, detail="Last message must be from user")

    username, user_base = require_user(request)
    file_rel = req.file.replace("\\", "/").strip("/")
    target = safe_join(user_base, file_rel)
    def ndjson_iter():
        # Prepare file context via MCP (and stream tool status)
        used_files: list[str] = []
        mcp_path = "/".join(part for part in [username, file_rel] if part)
        if target.exists() and target.is_file():
            # Start read_file tool
            try:
                yield (json.dumps({"type": "tool", "tool": "read_file", "status": "start", "path": mcp_path}, ensure_ascii=False) + "\n").encode("utf-8")
            except Exception:
                pass
            try:
                p, file_content = _read_single_file_for_chat(username, file_rel, request)
                used_files = [p]
                # success
                try:
                    yield (json.dumps({"type": "tool", "tool": "read_file", "status": "success", "path": p}, ensure_ascii=False) + "\n").encode("utf-8")
                except Exception:
                    pass
            except Exception as e:
                try:
                    yield (json.dumps({"type": "tool", "tool": "read_file", "status": "error", "path": mcp_path, "message": f"{e}"}, ensure_ascii=False) + "\n").encode("utf-8")
                except Exception:
                    pass
                file_content = ""
                used_files = []
        else:
            file_content = ""
            used_files = []

        system_prompt = _load_system_prompt()
        messages: list[dict] = [{"role": "system", "content": system_prompt}]
        messages.extend({"role": m.role, "content": m.content} for m in req.messages)
        messages.append(
            {
                "role": "system",
                "content": "The following is the content of the user's selected file. Use it to answer.\n\n"
                + f"FILE: {mcp_path}\n{file_content}",
            }
        )

        # Emit meta before tool loop
        yield (json.dumps({"type": "meta", "usedFiles": used_files}, ensure_ascii=False) + "\n").encode("utf-8")

        profile = _load_llm_profile(username)

        # If Agents SDK is available and OpenAI is preferred, run via Agents SDK
        if Agent is not None and MCPServerStreamableHttp is not None and profile.preferred == "openai":
            try:
                cfg = profile.openai
                if not cfg.apiKey:
                    raise RuntimeError("OpenAI API key is not configured")
                import asyncio
                async def _run() -> str:
                    prev_key = os.environ.get("OPENAI_API_KEY")
                    prev_base = os.environ.get("OPENAI_BASE_URL")
                    try:
                        os.environ["OPENAI_API_KEY"] = cfg.apiKey
                        if cfg.baseUrl:
                            os.environ["OPENAI_BASE_URL"] = cfg.baseUrl
                        model_name = (cfg.model or os.environ.get("OPENAI_MODEL") or "gpt-4o-mini")
                        async with MCPServerStreamableHttp(
                            name="Directory RAG MCP",
                            params={
                                "url": f"{MCP_SERVER_BASE}/mcp",
                                "headers": {"authorization": _get_auth_header(request)},
                                "timeout": 15,
                            },
                            cache_tools_list=True,
                        ) as server:
                            system_texts = [m["content"] for m in messages if m.get("role") == "system"]
                            instructions = "\n\n".join(system_texts) if system_texts else _load_system_prompt()
                            user_messages = [m for m in messages if m.get("role") != "system"]
                            agent = Agent(
                                name="Assistant",
                                instructions=instructions,
                                mcp_servers=[server],
                                model=model_name,
                                model_settings=ModelSettings(tool_choice="auto"),
                            )
                            transcript = []
                            for m in user_messages:
                                role = m.get("role") or "user"
                                content = m.get("content") or ""
                                prefix = "User:" if role == "user" else "Assistant:"
                                transcript.append(f"{prefix} {content}")
                            prompt = "\n".join(transcript)
                            result = await Runner.run(agent, prompt)
                            return result.final_output or ""
                    finally:
                        if prev_key is not None:
                            os.environ["OPENAI_API_KEY"] = prev_key
                        else:
                            os.environ.pop("OPENAI_API_KEY", None)
                        if prev_base is not None:
                            os.environ["OPENAI_BASE_URL"] = prev_base
                        else:
                            os.environ.pop("OPENAI_BASE_URL", None)
                answer = asyncio.run(_run())
                if answer:
                    yield (json.dumps({"type": "delta", "content": answer}, ensure_ascii=False) + "\n").encode("utf-8")
                yield (json.dumps({"type": "done"}, ensure_ascii=False) + "\n").encode("utf-8")
                return
            except Exception as e:
                try:
                    yield (json.dumps({"type": "error", "message": f"Agents SDK error: {e}"}, ensure_ascii=False) + "\n").encode("utf-8")
                except Exception:
                    pass
                # fall through to legacy tool loop

        # Directive completion helper
        def _complete_directive(msgs: list[dict]) -> tuple[str, dict]:
            try:
                if profile.preferred == "azure":
                    cfg = profile.azure
                    if not (cfg.apiKey and cfg.endpoint and cfg.deployment):
                        raise RuntimeError("Azure OpenAI is not fully configured")
                    client = AzureOpenAI(
                        azure_endpoint=cfg.endpoint,
                        api_key=cfg.apiKey,
                        api_version=(cfg.apiVersion or "2024-02-15-preview"),
                    )
                    comp = client.chat.completions.create(model=cfg.deployment, messages=msgs, temperature=0.1)
                else:
                    cfg = profile.openai
                    if not cfg.apiKey:
                        raise RuntimeError("OpenAI API key is not configured")
                    client = OpenAI(api_key=cfg.apiKey, base_url=(cfg.baseUrl or None))
                    model_name = (cfg.model or os.environ.get("OPENAI_MODEL") or "gpt-4o-mini")
                    comp = client.chat.completions.create(model=model_name, messages=msgs, temperature=0.1)
                text = (comp.choices[0].message.content or "").strip()
                first_line = text.splitlines()[0] if "\n" in text else text
                obj = json.loads(first_line)
                if isinstance(obj, dict) and obj.get("action") in ("tool", "final"):
                    return "ok", obj
                return "parse_error", {"raw": text}
            except Exception as ex:
                return "error", {"message": f"{ex}"}

        # Add tool protocol instructions to messages
        messages.append({"role": "system", "content": _tool_protocol_instructions()})

        # Interactive tool loop
        max_steps = 8
        step = 0
        last_answer = ""
        while step < max_steps:
            step += 1
            status, obj = _complete_directive(messages + [{"role": "system", "content": "Respond with a single JSON object as instructed."}])
            if status == "ok" and isinstance(obj, dict):
                if obj.get("action") == "final":
                    last_answer = str(obj.get("content") or "")
                    break
                if obj.get("action") == "tool":
                    tool = str(obj.get("tool") or "").strip()
                    args = obj.get("args") or {}
                    try:
                        yield (json.dumps({"type": "tool", "tool": tool, "status": "start", "path": str(args.get("path") or args.get("parent") or "")}, ensure_ascii=False) + "\n").encode("utf-8")
                    except Exception:
                        pass
                    try:
                        auth = _get_auth_header(request)
                        result: dict | list | str
                        if tool == "get_meta":
                            p = str(args.get("path") or "")
                            url = f"{MCP_SERVER_BASE}/tools/get_meta?" + urllib.parse.urlencode({"path": p})
                            result = _http_get_json(url, auth)
                        elif tool == "read_file":
                            p = str(args.get("path") or "")
                            url = f"{MCP_SERVER_BASE}/tools/read_file?" + urllib.parse.urlencode({"path": p})
                            result = _http_get_json(url, auth)
                            if isinstance(p, str) and p:
                                try:
                                    used_files.append(p)
                                except Exception:
                                    pass
                        elif tool == "create_file":
                            parent = str(args.get("parent") or "")
                            name = str(args.get("name") or "")
                            content = str(args.get("content") or "")
                            encoding = str(args.get("encoding") or "utf-8")
                            url = f"{MCP_SERVER_BASE}/tools/create_file"
                            result = _http_post_json(url, {"parent": parent, "name": name, "content": content, "encoding": encoding}, auth)
                        elif tool == "update_file":
                            p = str(args.get("path") or "")
                            content = str(args.get("content") or "")
                            encoding = str(args.get("encoding") or "utf-8")
                            url = f"{MCP_SERVER_BASE}/tools/update_file"
                            result = _http_post_json(url, {"path": p, "content": content, "encoding": encoding}, auth)
                        elif tool == "delete_file":
                            p = str(args.get("path") or "")
                            url = f"{MCP_SERVER_BASE}/tools/delete_file?" + urllib.parse.urlencode({"path": p})
                            result = _http_delete_json(url, auth)
                        else:
                            result = {"error": f"Unknown tool: {tool}"}
                        try:
                            created_path = str(args.get("path") or args.get("parent") or "")
                            if tool == "create_file":
                                try:
                                    nm = str(args.get("name") or "").strip()
                                    if created_path and nm:
                                        created_path = "/".join(p for p in [created_path.strip("/"), nm] if p)
                                except Exception:
                                    pass
                            yield (json.dumps({"type": "tool", "tool": tool, "status": "success", "path": created_path}, ensure_ascii=False) + "\n").encode("utf-8")
                        except Exception:
                            pass
                        try:
                            snippet = json.dumps(result, ensure_ascii=False) if not isinstance(result, str) else result
                        except Exception:
                            snippet = str(result)
                        messages.append({"role": "system", "content": f"Tool {tool} result (JSON):\n{snippet}"})
                        continue
                    except Exception as tex:
                        try:
                            yield (json.dumps({"type": "tool", "tool": tool, "status": "error", "message": f"{tex}"}, ensure_ascii=False) + "\n").encode("utf-8")
                        except Exception:
                            pass
                        messages.append({"role": "system", "content": f"Tool {tool} error: {tex}"})
                        continue
            break

        if not last_answer:
            try:
                if profile.preferred == "azure":
                    cfg = profile.azure
                    if not (cfg.apiKey and cfg.endpoint and cfg.deployment):
                        raise RuntimeError("Azure OpenAI is not fully configured")
                    client = AzureOpenAI(
                        azure_endpoint=cfg.endpoint,
                        api_key=cfg.apiKey,
                        api_version=(cfg.apiVersion or "2024-02-15-preview"),
                    )
                    stream = client.chat.completions.create(model=cfg.deployment, messages=messages, temperature=0.2, stream=True)
                else:
                    cfg = profile.openai
                    if not cfg.apiKey:
                        raise RuntimeError("OpenAI API key is not configured")
                    client = OpenAI(api_key=cfg.apiKey, base_url=(cfg.baseUrl or None))
                    model_name = (cfg.model or os.environ.get("OPENAI_MODEL") or "gpt-4o-mini")
                    stream = client.chat.completions.create(model=model_name, messages=messages, temperature=0.2, stream=True)
                for chunk in stream:  # type: ignore[assignment]
                    try:
                        choice = (chunk.choices[0] if getattr(chunk, "choices", None) else None)
                        delta_obj = getattr(choice, "delta", None) if choice is not None else None
                        piece = getattr(delta_obj, "content", None) if delta_obj is not None else None
                        if isinstance(piece, str) and piece:
                            yield (json.dumps({"type": "delta", "content": piece}, ensure_ascii=False) + "\n").encode("utf-8")
                    except Exception:
                        pass
                yield (json.dumps({"type": "done"}, ensure_ascii=False) + "\n").encode("utf-8")
                return
            except Exception as e:
                yield (json.dumps({"type": "error", "message": f"OpenAI error: {e}"}, ensure_ascii=False) + "\n").encode("utf-8")
                yield (json.dumps({"type": "done"}, ensure_ascii=False) + "\n").encode("utf-8")
                return

        if last_answer:
            yield (json.dumps({"type": "delta", "content": last_answer}, ensure_ascii=False) + "\n").encode("utf-8")
        yield (json.dumps({"type": "done"}, ensure_ascii=False) + "\n").encode("utf-8")

    return StreamingResponse(ndjson_iter(), media_type="application/x-ndjson")





class ShareCreateRequest(BaseModel):
    path: str
    expiresSec: int = 24 * 60 * 60  # default 24h


class ShareCreateResponse(BaseModel):
    url: str
    expires: int  # epoch seconds

class ShareStatusResponse(BaseModel):
    active: bool
    url: str | None = None
    expires: int | None = None

class ShareToggleRequest(BaseModel):
    path: str
    active: bool


@app.post("/api/share/create", response_model=ShareCreateResponse)
def create_share_link(request: Request, req: ShareCreateRequest):
    """Create a signed, time-limited share URL for a file.

    Returns a relative URL that can be appended to the server origin.
    """
    if req.expiresSec <= 0 or req.expiresSec > 30 * 24 * 60 * 60:
        # up to 30 days
        raise HTTPException(status_code=400, detail="Invalid expiresSec")

    # Only allow sharing files under the caller's own user directory
    _, user_base = require_user(request)
    file_path = safe_join(user_base, req.path)
    if not file_path.exists() or not file_path.is_file():
        raise HTTPException(status_code=404, detail="File not found")

    exp = int(time.time()) + int(req.expiresSec)
    # Encode the share path relative to the global BASE_DIR for download endpoint
    rel_from_base = str(file_path.relative_to(BASE_DIR)).replace(os.sep, "/")
    token = sign_share_token(rel_from_base, exp)
    # Return relative URL; frontend will prepend API base
    url = f"/api/share/download?path={urllib.parse.quote(rel_from_base)}&exp={exp}&token={token}"

    # Persist share status: mark this file as shared with latest expiry
    try:
        shares = _load_shares()
        # Store per-user namespace to avoid collisions between users with same filenames
        username = str(require_user(request)[0])
        user_map = shares.get(username) or {}
        # Keyed by user-relative path (req.path)
        user_map[req.path] = {"active": True, "url": url, "expires": exp}
        shares[username] = user_map
        _save_shares(shares)
    except Exception:
        # Non-fatal: link still returned even if persisting state failed
        pass

    return ShareCreateResponse(url=url, expires=exp)


@app.get("/api/share/download")
def share_download(
    path: str = Query(..., description="Relative file path under base directory"),
    exp: int = Query(..., description="Epoch seconds when the link expires"),
    token: str = Query(..., description="Signature token"),
):
    """Download a file via a signed share link."""
    # Fast path: check expiry
    now = int(time.time())
    if now > exp:
        raise HTTPException(status_code=410, detail="Link expired")

    # Validate signature then file path
    if not verify_share_token(path, exp, token):
        raise HTTPException(status_code=403, detail="Invalid token")

    # Enforce share registry OFF: path is relative to BASE_DIR
    try:
        # Determine owner and user-relative path
        # Expected shape: "{owner}/{user_rel}" because create_share_link used relative_to(BASE_DIR)
        owner, user_rel = path.split("/", 1)
        shares = _load_shares()
        user_map = shares.get(owner) or {}
        meta = user_map.get(user_rel)
        if meta and not bool(meta.get("active", False)):
            raise HTTPException(status_code=403, detail="Sharing disabled for this file")
    except ValueError:
        # If path doesn't include owner segment, skip registry check
        pass
    except HTTPException:
        raise
    except Exception:
        # On any error, default to allowing (signature already validated)
        pass

    file_path = safe_join(BASE_DIR, path)
    if not file_path.exists() or not file_path.is_file():
        raise HTTPException(status_code=404, detail="File not found")

    media_type, _ = mimetypes.guess_type(str(file_path))
    return FileResponse(
        str(file_path),
        filename=file_path.name,
        media_type=media_type or "application/octet-stream",
    )


@app.get("/api/share/status", response_model=ShareStatusResponse)
def share_status(request: Request, path: str = Query(..., description="Relative file path under user's base directory")):
    """Return whether the given file is currently shared and its link if active."""
    username, user_base = require_user(request)
    file_path = safe_join(user_base, path)
    if not file_path.exists() or not file_path.is_file():
        raise HTTPException(status_code=404, detail="File not found")
    shares = _load_shares()
    meta = (shares.get(username) or {}).get(path)
    if meta and bool(meta.get("active", False)):
        return ShareStatusResponse(active=True, url=str(meta.get("url") or None), expires=int(meta.get("expires") or 0) or None)
    return ShareStatusResponse(active=False, url=None, expires=None)


@app.post("/api/share/toggle", response_model=ShareStatusResponse)
def share_toggle(request: Request, body: ShareToggleRequest):
    """Enable or disable sharing for a specific file for the current user."""
    username, user_base = require_user(request)
    file_path = safe_join(user_base, body.path)
    if not file_path.exists() or not file_path.is_file():
        raise HTTPException(status_code=404, detail="File not found")
    shares = _load_shares()
    user_map = shares.get(username) or {}
    meta = user_map.get(body.path) or {}
    active = bool(body.active)
    meta["active"] = active
    # Keep existing url/expires if any
    user_map[body.path] = meta
    shares[username] = user_map
    _save_shares(shares)
    return ShareStatusResponse(active=active, url=str(meta.get("url") or None), expires=int(meta.get("expires") or 0) or None)


# ===================== Public Agent APIs =====================

class AgentCreateRequest(BaseModel):
    label: str | None = None
    # scope for the agent: either a folder or a single file
    folder: str | None = None
    file: str | None = None


class AgentCreateResponse(BaseModel):
    agentId: str
    label: str
    baseUrl: str  # base url for OpenAI-compatible API for this agent


def _validate_scope(user_base: Path, folder: str | None, file: str | None) -> tuple[str, str]:
    if (folder and file) or (not folder and not file):
        raise HTTPException(status_code=400, detail="Specify exactly one of folder or file")
    if folder:
        folder_rel = folder.replace("\\", "/").strip("/")
        target = safe_join(user_base, folder_rel)
        if not target.exists() or not target.is_dir():
            raise HTTPException(status_code=404, detail="Target folder not found")
        return ("folder", folder_rel)
    else:
        file_rel = (file or "").replace("\\", "/").strip("/")
        target = safe_join(user_base, file_rel)
        if not target.exists() or not target.is_file():
            raise HTTPException(status_code=404, detail="Target file not found")
        return ("file", file_rel)


@app.post("/api/agent/create", response_model=AgentCreateResponse)
def create_agent(request: Request, req: AgentCreateRequest):
    username, user_base = require_user(request)
    scope_kind, scope_path = _validate_scope(user_base, req.folder, req.file)

    # Generate agent id
    agent_id = secrets.token_urlsafe(12)
    label = req.label or (f"{scope_kind}:{scope_path}" if scope_path else scope_kind)
    agents = _load_agents()
    agents[agent_id] = {
        "owner": username,
        "kind": scope_kind,
        "path": scope_path,
        "label": label,
        "created": int(time.time()),
    }
    _save_agents(agents)

    # Public OpenAI-compatible base URL for this agent
    base_url = f"/api/public/agents/{agent_id}"
    return AgentCreateResponse(agentId=agent_id, label=label, baseUrl=base_url)


class AgentByScopeResponse(BaseModel):
    agentId: str
    label: str
    baseUrl: str


@app.get("/api/agent/by-scope", response_model=AgentByScopeResponse)
def get_agent_by_scope(request: Request, folder: str | None = Query(None), file: str | None = Query(None)):
    username, user_base = require_user(request)
    scope_kind, scope_path = _validate_scope(user_base, folder, file)
    found = _find_agent_by_scope(username, scope_kind, scope_path)
    if not found:
        raise HTTPException(status_code=404, detail="Agent not found")
    agent_id, meta = found
    base_url = f"/api/public/agents/{agent_id}"
    return AgentByScopeResponse(agentId=agent_id, label=str(meta.get("label") or ""), baseUrl=base_url)


class AgentDeleteResponse(BaseModel):
    ok: bool


@app.delete("/api/agent/by-scope", response_model=AgentDeleteResponse)
def delete_agent_by_scope(request: Request, folder: str | None = Query(None), file: str | None = Query(None)):
    username, user_base = require_user(request)
    scope_kind, scope_path = _validate_scope(user_base, folder, file)
    found = _find_agent_by_scope(username, scope_kind, scope_path)
    if not found:
        return AgentDeleteResponse(ok=True)
    agent_id, _ = found
    agents = _load_agents()
    if agent_id in agents:
        agents.pop(agent_id, None)
        _save_agents(agents)
    return AgentDeleteResponse(ok=True)


def _get_agent(agent_id: str) -> dict:
    agents = _load_agents()
    data = agents.get(agent_id)
    if not data:
        raise HTTPException(status_code=404, detail="Agent not found")
    # Validate the owner's directory still exists
    owner = str(data.get("owner") or "").strip().lower()
    if not owner:
        raise HTTPException(status_code=404, detail="Agent not found")
    return data


def _find_agent_by_scope(owner: str, scope_kind: str, scope_path: str) -> tuple[str, dict] | None:
    agents = _load_agents()
    for aid, meta in agents.items():
        try:
            if str(meta.get("owner")).lower() != owner.lower():
                continue
            if str(meta.get("kind")) != scope_kind:
                continue
            if str(meta.get("path")) != scope_path:
                continue
            return aid, meta
        except Exception:
            continue
    return None


def _agent_build_messages(owner: str, scope_kind: str, scope_path: str, user_msg: list[AgentMessage]) -> tuple[list[dict], list[str]]:
    # Build MCP context and collect files based on current logic
    mcp_folder = f"{owner}/{scope_path}" if scope_kind == "folder" and scope_path else owner
    used_files: list[str] = []
    file_chunks: list[str] = []
    system_prompt = _load_system_prompt()
    messages: list[dict] = [{"role": "system", "content": system_prompt}]
    for m in user_msg:
        messages.append({"role": m.role, "content": m.content})

    # Use internal MCP REST shims with owner token; issue a short-lived internal token for the owner
    users = _load_users()
    key = owner if owner in users else next((k for k in users.keys() if k.lower() == owner), owner)
    token_version = int((users.get(key) or {}).get("tokenVersion", 1))
    internal_token = _token_sign({"sub": owner, "v": token_version, "exp": int(time.time()) + 300})
    authz = f"Bearer {internal_token}"

    if scope_kind == "folder":
        meta_url = f"{MCP_SERVER_BASE}/tools/get_meta?{urllib.parse.urlencode({'path': mcp_folder})}"
        tree = _http_get_json(meta_url, authz)
        all_files = _flatten_tree(mcp_folder, tree if isinstance(tree, dict) else {}, "", True)
        question = user_msg[-1].content if user_msg else ""
        candidates = _select_candidate_files(question, all_files, 5)
        for p in candidates:
            rf_url = f"{MCP_SERVER_BASE}/tools/read_file?{urllib.parse.urlencode({'path': p})}"
            data = _http_get_json(rf_url, authz)
            content: str
            if isinstance(data, dict) and "content" in data and isinstance(data["content"], str):
                content = data["content"]
            else:
                content = json.dumps(data, ensure_ascii=False)
            if len(content) > 20000:
                content = content[:20000] + "\n...[truncated]"
            used_files.append(p)
            file_chunks.append(f"FILE: {p}\n" + content)
        if file_chunks:
            messages.append({"role": "system", "content": "The following are relevant files from the user's selected folder. Use them to answer.\n\n" + "\n\n".join(file_chunks)})
    else:
        # file
        mcp_path = "/".join(part for part in [owner, scope_path] if part)
        rf_url = f"{MCP_SERVER_BASE}/tools/read_file?{urllib.parse.urlencode({'path': mcp_path})}"
        data = _http_get_json(rf_url, authz)
        content = data.get("content") if isinstance(data, dict) else None  # type: ignore[assignment]
        if not isinstance(content, str):
            content = json.dumps(data, ensure_ascii=False)
        if len(content) > 20000:
            content = content[:20000] + "\n...[truncated]"
        used_files.append(mcp_path)
        messages.append({"role": "system", "content": "The following is the content of the user's selected file. Use it to answer.\n\n" + f"FILE: {mcp_path}\n{content}"})
    return messages, used_files


class PublicChatRequest(BaseModel):
    messages: list[AgentMessage]


class PublicChatResponse(BaseModel):
    messages: list[AgentMessage]
    usedFiles: list[str]


@app.post("/api/public/agents/{agent_id}/chat", response_model=PublicChatResponse)
def public_agent_chat(agent_id: str, req: PublicChatRequest):
    if not req.messages or req.messages[-1].role != "user":
        raise HTTPException(status_code=400, detail="Last message must be from user")
    agent = _get_agent(agent_id)
    owner = str(agent.get("owner"))
    scope_kind = str(agent.get("kind"))
    scope_path = str(agent.get("path"))
    # Build tool-first messages (do not prefetch files)
    system_prompt = _load_system_prompt()
    messages: list[dict] = [{"role": "system", "content": system_prompt}]
    messages.extend({"role": m.role, "content": m.content} for m in req.messages)
    messages.append({"role": "system", "content": _tool_protocol_instructions()})
    # Provide a light hint on scope to the model
    scope_label = f"{scope_kind}:{scope_path}" if scope_path else scope_kind
    owner_hint = str(owner)
    messages.append({"role": "system", "content": f"Agent scope: {scope_label}. Owner: {owner_hint}. Use absolute MCP paths like '{owner_hint}/...'."})

    # Use owner's profile and an internal auth header for MCP
    profile = _load_llm_profile(owner)
    users = _load_users()
    key = owner if owner in users else next((k for k in users.keys() if k.lower() == owner), owner)
    token_version = int((users.get(key) or {}).get("tokenVersion", 1))
    internal_token = _token_sign({"sub": owner, "v": token_version, "exp": int(time.time()) + 300})
    authz = f"Bearer {internal_token}"

    answer, used_files = _run_tool_calling(messages, profile, authz)
    out_messages = list(req.messages)
    out_messages.append(AgentMessage(role="assistant", content=answer or ""))
    return PublicChatResponse(messages=out_messages, usedFiles=used_files)


# NDJSON streaming version for public chat (no auth required)
@app.post("/api/public/agents/{agent_id}/chat/stream")
def public_agent_chat_stream(agent_id: str, req: PublicChatRequest):
    if not req.messages or req.messages[-1].role != "user":
        raise HTTPException(status_code=400, detail="Last message must be from user")
    agent = _get_agent(agent_id)
    owner = str(agent.get("owner"))
    scope_kind = str(agent.get("kind"))
    scope_path = str(agent.get("path"))

    def ndjson_iter():
        # Build messages similarly to non-streaming public handler
        messages, used_files = _agent_build_messages(owner, scope_kind, scope_path, req.messages)
        # Emit meta about used files (if any were preselected)
        try:
            yield (json.dumps({"type": "meta", "usedFiles": used_files}, ensure_ascii=False) + "\n").encode("utf-8")
        except Exception:
            pass

        profile = _load_llm_profile(owner)
        # Try streaming directly from provider
        try:
            if profile.preferred == "azure":
                cfg = profile.azure
                if not (cfg.apiKey and cfg.endpoint and cfg.deployment):
                    raise RuntimeError("Azure OpenAI is not fully configured for agent owner")
                client = AzureOpenAI(azure_endpoint=cfg.endpoint, api_key=cfg.apiKey, api_version=(cfg.apiVersion or "2024-02-15-preview"))
                stream = client.chat.completions.create(model=cfg.deployment, messages=messages, temperature=0.2, stream=True)
            else:
                cfg = profile.openai
                if not cfg.apiKey:
                    raise RuntimeError("OpenAI API key is not configured for agent owner")
                client = OpenAI(api_key=cfg.apiKey, base_url=(cfg.baseUrl or None))
                model_name = (cfg.model or os.environ.get("OPENAI_MODEL") or "gpt-4o-mini")
                stream = client.chat.completions.create(model=model_name, messages=messages, temperature=0.2, stream=True)
            for chunk in stream:  # type: ignore[assignment]
                try:
                    choice = (chunk.choices[0] if getattr(chunk, "choices", None) else None)
                    delta_obj = getattr(choice, "delta", None) if choice is not None else None
                    piece = getattr(delta_obj, "content", None) if delta_obj is not None else None
                    if isinstance(piece, str) and piece:
                        yield (json.dumps({"type": "delta", "content": piece}, ensure_ascii=False) + "\n").encode("utf-8")
                except Exception:
                    pass
            yield (json.dumps({"type": "done"}, ensure_ascii=False) + "\n").encode("utf-8")
            return
        except Exception as e:
            try:
                yield (json.dumps({"type": "error", "message": f"LLM error: {e}"}, ensure_ascii=False) + "\n").encode("utf-8")
            except Exception:
                pass
            yield (json.dumps({"type": "done"}, ensure_ascii=False) + "\n").encode("utf-8")
            return

    return StreamingResponse(ndjson_iter(), media_type="application/x-ndjson")

# OpenAI-compatible Completions endpoint for this agent
@app.post("/api/public/agents/{agent_id}/v1/chat/completions")
def public_agent_chat_completions(agent_id: str, body: dict):
    # Minimal OpenAI-compatible handler: expects { model, messages, ... }
    # We ignore model; we route based on owner's profile
    messages_in = body.get("messages")
    if not isinstance(messages_in, list) or not messages_in:
        raise HTTPException(status_code=400, detail="messages is required")
    if not (isinstance(messages_in[-1], dict) and messages_in[-1].get("role") == "user"):
        raise HTTPException(status_code=400, detail="Last message must be from user")
    agent = _get_agent(agent_id)
    owner = str(agent.get("owner"))
    scope_kind = str(agent.get("kind"))
    scope_path = str(agent.get("path"))

    # Coerce incoming messages to AgentMessage list
    user_msgs: list[AgentMessage] = []
    for m in messages_in:
        role = str(m.get("role") or "")
        content = m.get("content")
        if isinstance(content, list):
            # If content is list (tool messages unsupported here), join text parts
            parts = []
            for c in content:
                if isinstance(c, dict) and c.get("type") == "text" and isinstance(c.get("text"), str):
                    parts.append(c["text"])  # type: ignore[index]
            content = "\n".join(parts)
        if not isinstance(content, str):
            content = ""
        user_msgs.append(AgentMessage(role=role, content=content))

    messages, used_files = _agent_build_messages(owner, scope_kind, scope_path, user_msgs)

    profile = _load_llm_profile(owner)
    preferred = profile.preferred
    answer = ""
    last_error: Exception | None = None

    def _try_openai() -> str:
        cfg = profile.openai
        if not cfg.apiKey:
            raise RuntimeError("OpenAI API key is not configured for agent owner")
        client = OpenAI(api_key=cfg.apiKey, base_url=(cfg.baseUrl or None))
        model_name = (cfg.model or os.environ.get("OPENAI_MODEL") or "gpt-4o-mini")
        comp = client.chat.completions.create(model=model_name, messages=messages, temperature=0.2)  # type: ignore[arg-type]
        return comp.choices[0].message.content if comp.choices else ""

    def _try_azure() -> str:
        cfg = profile.azure
        if not (cfg.apiKey and cfg.endpoint and cfg.deployment):
            raise RuntimeError("Azure OpenAI is not fully configured for agent owner")
        client = AzureOpenAI(azure_endpoint=cfg.endpoint, api_key=cfg.apiKey, api_version=(cfg.apiVersion or "2024-02-15-preview"))
        comp = client.chat.completions.create(model=cfg.deployment, messages=messages, temperature=0.2)  # type: ignore[arg-type]
        return comp.choices[0].message.content if comp.choices else ""

    for p in [preferred, ("azure" if preferred == "openai" else "openai")]:
        try:
            answer = _try_azure() if p == "azure" else _try_openai()
            last_error = None
            break
        except Exception as e:
            last_error = e
            continue
    if last_error is not None and not answer:
        raise HTTPException(status_code=502, detail=f"LLM error: {last_error}")

    # Return OpenAI-compatible response shape
    return {
        "id": f"chatcmpl_{secrets.token_hex(8)}",
        "object": "chat.completion",
        "created": int(time.time()),
        "model": body.get("model") or "",
        "choices": [
            {
                "index": 0,
                "message": { "role": "assistant", "content": answer },
                "finish_reason": "stop",
            }
        ],
        "usage": { "prompt_tokens": 0, "completion_tokens": 0, "total_tokens": 0 },
    }


# Dev-only graceful shutdown endpoint to release :9000 when closing public-agent
@app.post("/api/dev/shutdown")
def dev_shutdown():
    flag = os.environ.get("ALLOW_DEV_SHUTDOWN", "0").lower()
    if flag not in {"1", "true", "yes"}:
        return {"ok": False, "message": "disabled"}
    def _exit_later():
        try:
            time.sleep(0.2)
        finally:
            os._exit(0)
    try:
        threading.Thread(target=_exit_later, daemon=True).start()
    except Exception:
        # best-effort
        os._exit(0)
    return {"ok": True}
